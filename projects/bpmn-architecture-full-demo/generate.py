#!/usr/bin/env python3
"""BPMN Architecture full demo generator (data-driven + dynamic loaders).

Orchestrator — reads config YAML, discovers BP loaders, produces PPTX.
"""

import importlib.util
import re
import sys
from pathlib import Path

import yaml
from pptx.dml.color import RGBColor
from pptx.util import Pt

ROOT = Path(__file__).resolve().parent.parent.parent
SCRIPTS_DIR = ROOT / "skills" / "pptx" / "scripts"
LOADERS_DIR = ROOT / "skills" / "pptx" / "bpmn-architecture" / "loaders"
SHARED_DIR = ROOT / "skills" / "pptx" / "shared"

sys.path.insert(0, str(SCRIPTS_DIR))
sys.path.insert(0, str(SHARED_DIR))
from pptx_lib import build_pptx, header, layout_by_names, verify_pptx  # noqa: E402
from renderer_utils import configure_theme  # noqa: E402

TEMPLATE = ROOT / "skills" / "pptx" / "master-library" / "dark-cloudwise-green" / "cloudwise-master.pptx"
TEMPLATE_SPEC = ROOT / "skills" / "pptx" / "master-library" / "dark-cloudwise-green" / "cloudwise-spec.yaml"
OUTPUT = Path(__file__).resolve().parent / "bpmn-architecture-full-demo.pptx"
CONFIG_FILE = Path(__file__).resolve().parent / "data" / "config_template.yaml"

LOADER_REGISTRY = {}


def load_config(path):
    with open(path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def register_loader(name, load_fn):
    LOADER_REGISTRY[name] = load_fn


def load_all_loaders():
    files = sorted(LOADERS_DIR.glob("bp_*.py"))
    print(f"[Info] 找到 {len(files)} 个 BP 装载器模块")
    for loader_file in files:
        module_name = loader_file.stem
        spec = importlib.util.spec_from_file_location(module_name, loader_file)
        module = importlib.util.module_from_spec(spec)
        try:
            spec.loader.exec_module(module)
            if hasattr(module, "load_slide"):
                register_loader(module_name, module.load_slide)
                print(f"  ✓ {module_name}")
            else:
                print(f"  ✗ {module_name} (缺少 load_slide)")
        except Exception as exc:
            print(f"  ✗ {module_name} (加载失败: {exc})")


def sort_loader_names(names):
    def _key(name):
        m = re.match(r"bp_(\d+)_", name)
        return int(m.group(1)) if m else 999
    return sorted(names, key=_key)


def add_cover_slide(ctx, cfg):
    slide = ctx.add_cover_slide()
    title_ph = next(
        (s for s in slide.shapes
         if getattr(s, "is_placeholder", False) and s.placeholder_format.idx == 0),
        None,
    )
    if title_ph and title_ph.has_text_frame:
        tf = title_ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = cfg.get("title", "BPMN Architecture Demo")
        p.font.size = Pt(48)
        p.font.bold = True

    subtitle_ph = next(
        (s for s in slide.shapes
         if getattr(s, "is_placeholder", False) and s.placeholder_format.idx == 1),
        None,
    )
    if subtitle_ph and subtitle_ph.has_text_frame:
        tf = subtitle_ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = cfg.get("subtitle", "Data-driven BPMN Diagram Collection")
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xA5, 0xA7, 0xAA)


def add_chapter_slide(ctx, title_text):
    layout = layout_by_names(ctx.prs, ["节标题", "刻度线", "標題", "标题"], 0)
    if not layout:
        layout = layout_by_names(ctx.prs, ["标题", "標題"], 0)
    slide = ctx.prs.slides.add_slide(layout)
    title_ph = next(
        (s for s in slide.shapes
         if getattr(s, "is_placeholder", False) and s.placeholder_format.idx == 0),
        None,
    )
    if title_ph and title_ph.has_text_frame:
        tf = title_ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(42)
        p.font.bold = True


def add_back_slide(ctx, cfg):
    slide = ctx.add_thanks_slide()
    title_ph = next(
        (s for s in slide.shapes
         if getattr(s, "is_placeholder", False) and s.placeholder_format.idx == 0),
        None,
    )
    if title_ph and title_ph.has_text_frame:
        tf = title_ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = cfg.get("back_title", "Thank You")
        p.font.size = Pt(54)
        p.font.bold = True


def my_slides(ctx):
    configure_theme(ctx.colors, ctx.template_spec)
    cfg = load_config(CONFIG_FILE)

    add_cover_slide(ctx, cfg.get("meta", {}))

    sections = cfg.get("sections", [])
    for section in sections:
        section_title = section.get("title", "BPMN Section")
        slides = section.get("slides", {})
        add_chapter_slide(ctx, section_title)

        for loader_name in sort_loader_names(list(slides.keys())):
            loader_fn = LOADER_REGISTRY.get(loader_name)
            if not loader_fn:
                print(f"[Warn] {loader_name} 未注册")
                continue
            try:
                print(f"[Slide] 正在生成 {loader_name}...")
                loader_fn(ctx, slides[loader_name])
                print("  ✓ 成功")
            except Exception as exc:
                print(f"  ✗ 失败: {exc}")
                import traceback
                traceback.print_exc()

    add_back_slide(ctx, cfg.get("meta", {}))


if __name__ == "__main__":
    print("=" * 60)
    print("BPMN Architecture 全量生成器")
    print("=" * 60)

    print("\n[Step 1] 加载装载器...")
    load_all_loaders()
    print(f"已注册 {len(LOADER_REGISTRY)} 个装载器\n")

    print("[Step 2] 生成PPT...")
    build_pptx(TEMPLATE, OUTPUT, my_slides, TEMPLATE_SPEC)
    print(f"✓ PPT已生成: {OUTPUT}")

    print("[Step 3] 验证...")
    verify_pptx(OUTPUT, TEMPLATE)
    print("✓ 验证完成")
