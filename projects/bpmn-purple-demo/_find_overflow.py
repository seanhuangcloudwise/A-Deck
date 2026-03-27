"""Find which shapes overflow the slide x boundary."""
import sys
sys.path.insert(0, '/Volumes/work/Workspace/A-Deck/skills/pptx/scripts')
from pptx import Presentation
from pptx.util import Emu

prs = Presentation('/Volumes/work/Workspace/A-Deck/projects/bpmn-purple-demo/bpmn-purple-demo.pptx')
W = prs.slide_width  # 10" in EMU

for i, slide in enumerate(prs.slides):
    title = ''
    for s in slide.shapes:
        if getattr(s, 'is_placeholder', False) and s.placeholder_format.idx == 0:
            title = getattr(s, 'text', '')[:40]
    for shape in slide.shapes:
        right = shape.left + shape.width
        if right > W + Emu(9144):  # 0.01" tolerance
            right_in = right / 914400
            name = shape.name[:40]
            print(f'Slide {i+1} "{title}": shape "{name}" right={right_in:.3f}"  (overflow by {(right-W)/914400:.3f}")')
