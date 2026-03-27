#!/usr/bin/env python3
"""BPMN Architecture demo — light-cloudwise-purple master.

Uses the rich config YAML to exercise all 10 diagram types and their features.
"""

import importlib.util
import io
import re
import sys
import zipfile
from pathlib import Path

import yaml
from lxml import etree
from pptx.dml.color import RGBColor
from pptx.util import Pt

ROOT = Path(__file__).resolve().parent.parent.parent
SCRIPTS_DIR = ROOT / "skills" / "pptx" / "scripts"
LOADERS_DIR = ROOT / "skills" / "pptx" / "bpmn-architecture" / "loaders"
SHARED_DIR = ROOT / "skills" / "pptx" / "shared"

sys.path.insert(0, str(SCRIPTS_DIR))
sys.path.insert(0, str(SHARED_DIR))
from pptx_lib import build_pptx, layout_by_names, verify_pptx  # noqa: E402
from renderer_utils import configure_theme  # noqa: E402

TEMPLATE = ROOT / "skills" / "pptx" / "master-library" / "light-cloudwise-purple" / "cloudwise-master.pptx"
TEMPLATE_SPEC = ROOT / "skills" / "pptx" / "master-library" / "light-cloudwise-purple" / "cloudwise-spec.yaml"
OUTPUT = Path(__file__).resolve().parent / "bpmn-purple-demo.pptx"
CONFIG_FILE = Path(__file__).resolve().parent / "config.yaml"

LOADER_REGISTRY = {}


def load_config(path):
    with open(path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def register_loader(name, load_fn):
    LOADER_REGISTRY[name] = load_fn


def load_all_loaders():
    files = sorted(LOADERS_DIR.glob("bp_*.py"))
    print(f"[Info] 找到 {len(files)} 个 BP 装载器模块")
    for loader_file in files:
        module_name = loader_file.stem
        spec = importlib.util.spec_from_file_location(module_name, loader_file)
        module = importlib.util.module_from_spec(spec)
        try:
            spec.loader.exec_module(module)
            if hasattr(module, "load_slide"):
                register_loader(module_name, module.load_slide)
                print(f"  + {module_name}")
            else:
                print(f"  - {module_name} (no load_slide)")
        except Exception as exc:
            print(f"  ! {module_name}: {exc}")
            import traceback; traceback.print_exc()


def sort_loader_names(names):
    def _key(name):
        m = re.match(r"bp_(\d+)_", name)
        return int(m.group(1)) if m else 999
    return sorted(names, key=_key)


def add_cover_slide(ctx, cfg):
    slide = ctx.add_cover_slide()
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        idx = shape.placeholder_format.idx
        if idx == 0 and shape.has_text_frame:
            tf = shape.text_frame; tf.clear()
            p = tf.paragraphs[0]; p.text = cfg.get("title", "BPMN Demo")
            p.font.size = Pt(40); p.font.bold = True
        elif idx == 1 and shape.has_text_frame:
            tf = shape.text_frame; tf.clear()
            p = tf.paragraphs[0]; p.text = cfg.get("subtitle", "")
            p.font.size = Pt(14)


def add_chapter_slide(ctx, title_text):
    layout = layout_by_names(ctx.prs, ["节标题", "标题"], 0)
    slide = ctx.prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False) and shape.placeholder_format.idx == 0:
            if shape.has_text_frame:
                tf = shape.text_frame; tf.clear()
                p = tf.paragraphs[0]; p.text = title_text
                p.font.size = Pt(36); p.font.bold = True


def add_back_slide(ctx, cfg):
    slide = ctx.add_thanks_slide()
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False) and shape.placeholder_format.idx == 0:
            if shape.has_text_frame:
                tf = shape.text_frame; tf.clear()
                p = tf.paragraphs[0]; p.text = cfg.get("back_title", "Thank You")
                p.font.size = Pt(48); p.font.bold = True


def my_slides(ctx):
    configure_theme(ctx.colors, ctx.template_spec)
    cfg = load_config(CONFIG_FILE)

    add_cover_slide(ctx, cfg.get("meta", {}))

    for section in cfg.get("sections", []):
        add_chapter_slide(ctx, section.get("title", ""))
        for loader_name in sort_loader_names(list(section.get("slides", {}).keys())):
            loader_fn = LOADER_REGISTRY.get(loader_name)
            if not loader_fn:
                print(f"  [skip] {loader_name} not registered")
                continue
            try:
                print(f"  [slide] {loader_name}")
                loader_fn(ctx, section["slides"][loader_name])
                print("         OK")
            except Exception as exc:
                print(f"         FAIL: {exc}")
                import traceback; traceback.print_exc()

    add_back_slide(ctx, cfg.get("meta", {}))


# ---------------------------------------------------------------------------
# Post-processing: shadow removal
# ---------------------------------------------------------------------------

_DML_NS_SHADOW = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_SHADOW_TAGS = {
    f'{{{_DML_NS_SHADOW}}}outerShdw',
    f'{{{_DML_NS_SHADOW}}}innerShdw',
    f'{{{_DML_NS_SHADOW}}}prstShdw',
    f'{{{_DML_NS_SHADOW}}}reflection',
    f'{{{_DML_NS_SHADOW}}}glow',
    f'{{{_DML_NS_SHADOW}}}softEdge',
}


def _strip_shadows_from_pptx(pptx_path):
    """Strip all shadow/glow/reflection effects from theme XML in a PPTX file."""
    total = 0
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        names = zin.namelist()
        targets = [n for n in names if n.endswith('.xml')]
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for name in names:
                data = zin.read(name)
                if name in targets:
                    root = etree.fromstring(data)
                    removed = 0
                    for lst in root.iter(f'{{{_DML_NS_SHADOW}}}effectLst'):
                        for child in list(lst):
                            if child.tag in _SHADOW_TAGS:
                                lst.remove(child)
                                removed += 1
                    total += removed
                    new_data = etree.tostring(
                        root, xml_declaration=True,
                        encoding='UTF-8', standalone=True)
                    zout.writestr(name, new_data)
                else:
                    zout.writestr(name, data)

    with open(pptx_path, 'wb') as f:
        f.write(buf.getvalue())
    print(f"    已移除 {total} 个阴影/光晕元素")


if __name__ == "__main__":
    print("=" * 60)
    print("BPMN Purple Demo 生成器")
    print(f"母版: light-cloudwise-purple")
    print("=" * 60)

    print("\n[1] 加载装载器...")
    load_all_loaders()
    print(f"    已注册 {len(LOADER_REGISTRY)} 个装载器\n")

    print("[2] 生成 PPT...")
    build_pptx(TEMPLATE, OUTPUT, my_slides, TEMPLATE_SPEC)
    print(f"\n    输出: {OUTPUT}")

    print("\n[3] 移除阴影特效...")
    _strip_shadows_from_pptx(OUTPUT)

    print("\n[4] 验证...")
    verify_pptx(OUTPUT, TEMPLATE)
    print("    验证通过")


# ---------------------------------------------------------------------------
# Post-processing: shadow removal
# ---------------------------------------------------------------------------

_DML_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_SHADOW_TAGS = {
    f'{{{_DML_NS}}}outerShdw',
    f'{{{_DML_NS}}}innerShdw',
    f'{{{_DML_NS}}}prstShdw',
    f'{{{_DML_NS}}}reflection',
    f'{{{_DML_NS}}}glow',
    f'{{{_DML_NS}}}softEdge',
}


def _strip_shadows_from_pptx(pptx_path):
    """Strip all shadow/glow/reflection effects from theme XML in a PPTX file."""
    total = 0
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        names = zin.namelist()
        targets = [n for n in names if n.endswith('.xml')]
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for name in names:
                data = zin.read(name)
                if name in targets:
                    root = etree.fromstring(data)
                    removed = 0
                    for lst in root.iter(f'{{{_DML_NS}}}effectLst'):
                        for child in list(lst):
                            if child.tag in _SHADOW_TAGS:
                                lst.remove(child)
                                removed += 1
                    total += removed
                    new_data = etree.tostring(
                        root, xml_declaration=True,
                        encoding='UTF-8', standalone=True)
                    zout.writestr(name, new_data)
                else:
                    zout.writestr(name, data)

    with open(pptx_path, 'wb') as f:
        f.write(buf.getvalue())
    print(f"    已移除 {total} 个阴影/光晕元素")
