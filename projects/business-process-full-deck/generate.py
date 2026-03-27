#!/usr/bin/env python3
"""Generate A-Deck business-process deck (BA-03 layered) with Cloudwise master."""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent / "skills" / "pptx" / "scripts"))
from pptx_lib import *  # noqa: E402,F403

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TEMPLATE = (Path(__file__).resolve().parent.parent.parent
            / "skills" / "pptx" / "master-library"
            / "light-cloudwise-cyan" / "cloudwise-master.pptx")
OUTPUT = Path(__file__).resolve().parent / "a-deck-business-process.pptx"


def _i(v):
    return int(v)


def _lane(slide, x, y, w, h, title):
    rect(slide, x, y, Inches(1.22), h, C["domain_bg"])
    textbox(slide, x + Inches(0.08), y + Inches(0.08), Inches(1.06), h - Inches(0.16),
            title, size=10, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    rrect(slide, x + Inches(1.24), y, w - Inches(1.24), h, C["white"], line=C["line"], adj=2500)


def _step(slide, x, y, w, h, name, system="", duration="", critical=False, manual=False):
    fill = C["cyan_2"] if not manual else C["light"]
    edge = C["cyan"] if critical else C["line"]
    shape = rrect(slide, x, y, w, h, fill, line=edge, adj=2600)
    shape.line.width = Pt(2.0 if critical else 1.0)
    if manual:
        set_dash(shape)

    textbox(slide, x + Inches(0.06), y + Inches(0.08), w - Inches(0.12), Inches(0.24),
            name, size=9, color=C["dark"], align=PP_ALIGN.CENTER)

    if system:
        pill = rrect(slide, x + Inches(0.14), y + Inches(0.36), w - Inches(0.28), Inches(0.14),
                     C["domain_bg"], line=C["domain_bg"], adj=2200)
        add_text(pill, system, size=7, color=C["white"], align=PP_ALIGN.CENTER)

    if duration:
        dur = rrect(slide, x + w - Inches(0.34), y + Inches(0.02), Inches(0.30), Inches(0.12),
                    C["light"], line=C["brand_gray"], adj=1800)
        add_text(dur, duration, size=6, color=C["gray"], align=PP_ALIGN.CENTER)

    if critical:
        textbox(slide, x + w - Inches(0.10), y - Inches(0.03), Inches(0.08), Inches(0.08),
                "!", size=8, bold=True, color=RGBColor(0xC0, 0x39, 0x2B), align=PP_ALIGN.CENTER)


def _decision(slide, x, y, w, h, text):
    d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, _i(x), _i(y), _i(w), _i(h))
    d.fill.solid()
    d.fill.fore_color.rgb = C["white"]
    d.line.color.rgb = C["dark"]
    d.line.width = Pt(1.2)
    textbox(slide, x + Inches(0.04), y + Inches(0.13), w - Inches(0.08), Inches(0.2),
            text, size=8, color=C["dark"], align=PP_ALIGN.CENTER)


def _event(slide, x, y, start=True):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, _i(x), _i(y), _i(Inches(0.18)), _i(Inches(0.18)))
    c.fill.solid()
    c.fill.fore_color.rgb = C["domain_bg"] if start else C["dark"]
    c.line.color.rgb = C["domain_bg"] if start else C["dark"]


def _link(slide, x1, y1, x2, y2, label="", dashed=False, color=None, style="auto"):
    line_color = color or C["dark"]

    def _segment(sx, sy, ex, ey):
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _i(sx), _i(sy), _i(ex), _i(ey))
        conn.line.color.rgb = line_color
        conn.line.width = Pt(1.3)
        if dashed:
            conn.line.dash_style = 4
        return conn

    is_elbow = style == "elbow" or (style == "auto" and _i(y1) != _i(y2))
    if is_elbow:
        bend_x = x1 + (x2 - x1) * 0.5
        _segment(x1, y1, bend_x, y1)
        _segment(bend_x, y1, bend_x, y2)
        conn = _segment(bend_x, y2, x2, y2)
        if label:
            textbox(slide, bend_x + Inches(0.03), min(y1, y2) - Inches(0.12), Inches(0.9), Inches(0.12),
                    label, size=7, color=C["brand_gray"], align=PP_ALIGN.LEFT)
        return conn

    conn = _segment(x1, y1, x2, y2)
    if label:
        textbox(slide, min(x1, x2) + Inches(0.03), min(y1, y2) - Inches(0.12), Inches(0.9), Inches(0.12),
                label, size=7, color=C["brand_gray"], align=PP_ALIGN.LEFT)
    return conn


def _add_content_slide(ctx):
    """Force content pages to use layout named '内容' and keep placeholders."""
    layout = layout_by_names(ctx.prs, ["内容"], 1)
    return ctx.prs.slides.add_slide(layout)


def _subtitle(ctx, slide, text):
    subtitle_ph = next(
        (s for s in slide.shapes
         if getattr(s, "is_placeholder", False) and s.placeholder_format.idx == 1),
        None,
    )
    if subtitle_ph is not None and subtitle_ph.has_text_frame:
        tf = subtitle_ph.text_frame
        tf.clear()
        tf.paragraphs[0].text = text
        return
    textbox(slide, ctx.safe_left + Inches(0.04), Inches(0.68), ctx.safe_width - Inches(0.08), Inches(0.22),
            text, size=10, color=C["brand_gray"])


def _footer_meta(ctx, slide, text):
    textbox(slide, ctx.safe_left, ctx.safe_bottom + Inches(0.02), ctx.safe_width, Inches(0.18),
            text, size=7, color=C["brand_gray"], align=PP_ALIGN.LEFT)


def _step_positions(left, width, count, step_w, start_pad=1.36, end_pad=1.12, min_gap=0.10):
    body_left = left + Inches(start_pad)
    body_right = left + width - Inches(end_pad)
    if count <= 1:
        return [body_left]
    gap = (body_right - body_left - step_w * count) / (count - 1)
    gap = max(gap, Inches(min_gap))
    return [body_left + i * (step_w + gap) for i in range(count)]


def slide_l0(ctx):
    slide = _add_content_slide(ctx)
    header(slide, "PPT Maker Agent — Business Process Architecture")
    _subtitle(ctx, slide, "Layer 0 (Value Chain): End-to-End Process Landscape")

    phases = [
        ("P0", "需求受理", "需求入口与范围确认"),
        ("P1", "知识准备", "资料/母版/规范装载"),
        ("P2", "提纲与规划", "结构规划与页级策略"),
        ("P3", "内容与制图", "正文生成与架构图绘制"),
        ("P4", "校验与评审", "一致性与质量门禁"),
        ("P5", "发布与学习", "交付归档与经验回流"),
    ]

    x = ctx.safe_left + Inches(0.2)
    y = ctx.content_top + Inches(0.64)
    gap = Inches(0.12)
    w = (ctx.safe_width - Inches(0.4) - gap * 5) / 6
    h = Inches(1.05)

    for i, (pid, name, desc) in enumerate(phases):
        fill = C["cyan"] if i in (2, 3, 4) else C["cyan_2"]
        text_c = C["white"] if i in (2, 3, 4) else C["dark"]
        box = rrect(slide, x + i * (w + gap), y, w, h, fill, line=C["cyan"], adj=3200)
        box.line.width = Pt(1.2)
        textbox(slide, x + i * (w + gap) + Inches(0.06), y + Inches(0.10), w - Inches(0.12), Inches(0.24),
                f"{pid} · {name}", size=10, bold=True, color=text_c, align=PP_ALIGN.CENTER)
        textbox(slide, x + i * (w + gap) + Inches(0.08), y + Inches(0.42), w - Inches(0.16), Inches(0.5),
                desc, size=8, color=text_c, align=PP_ALIGN.CENTER)

        if i < len(phases) - 1:
            textbox(slide, x + i * (w + gap) + w + Inches(0.01), y + Inches(0.42), Inches(0.1), Inches(0.2),
                    "▶", size=12, color=C["cyan"], align=PP_ALIGN.CENTER)

    panel_top = y + h + Inches(0.41)
    panel_h = ctx.safe_bottom - panel_top - Inches(0.08)
    panel = rrect(slide, ctx.safe_left + Inches(0.22), panel_top, ctx.safe_width - Inches(0.44), panel_h,
                  C["light"], line=C["line"], adj=2200)
    panel.line.width = Pt(1.0)

    textbox(slide, ctx.safe_left + Inches(0.38), panel_top + Inches(0.16), Inches(2.8), Inches(0.24),
            "Layer Decomposition", size=11, bold=True, color=C["dark"])

    details = [
        "L1: Phase Process (4 pages) — Intake/Preparation/Generation/Delivery",
        "L2: Detailed Swimlane SOP (3 pages) — branch, exception, SLA",
        "Cross-layer Trace: BA-03 -> AA-03 integration points",
        "Owner: Product Ops  |  Version: v2.0  |  Baseline SLA: T+1 day",
    ]
    for i, t in enumerate(details):
        textbox(slide, ctx.safe_left + Inches(0.42), panel_top + Inches(0.54 + i * 0.44), ctx.safe_width - Inches(0.9), Inches(0.3),
                f"• {t}", size=9, color=C["gray"])

    _footer_meta(ctx, slide, "Scope: PPT Maker Agent end-to-end business process map · Cloudwise Master")


def slide_l1_intake(ctx):
    slide = _add_content_slide(ctx)
    header(slide, "L1 Process P0-P2 — Intake & Planning")
    _subtitle(ctx, slide, "Layer 1: 业务受理 -> 分析澄清 -> 提纲规划")

    left = ctx.safe_left
    top = ctx.content_top
    width = ctx.safe_width
    lane_h = Inches(1.47)

    _lane(slide, left, top + lane_h * 0, width, lane_h, "业务方")
    _lane(slide, left, top + lane_h * 1, width, lane_h, "产品经理")
    _lane(slide, left, top + lane_h * 2, width, lane_h, "架构师")

    sw = Inches(1.0)
    x = _step_positions(left, width, 6, sw)
    sh = Inches(0.58)
    y0 = top + Inches(0.44)
    y1 = top + lane_h + Inches(0.44)
    y2 = top + lane_h * 2 + Inches(0.44)

    _event(slide, x[0] - Inches(0.24), y0 + Inches(0.2), start=True)
    _step(slide, x[0], y0, sw, sh, "提交需求", "Portal", "<=2h")
    _step(slide, x[1], y1, sw, sh, "需求拆解", "Jira", "<=4h")
    _step(slide, x[2], y1, sw, sh, "范围确认", "Jira", "<=2h", critical=True)
    _step(slide, x[3], y2, sw, sh, "约束识别", "EA Repo", "<=3h")
    _step(slide, x[4], y1, sw, sh, "提纲规划", "Planner", "<=2h", critical=True)
    _step(slide, x[5], y1, sw, sh, "计划冻结", "Workflow", "<=1h")
    _event(slide, x[5] + Inches(1.12), y1 + Inches(0.2), start=False)

    _link(slide, x[0] + sw, y0 + Inches(0.29), x[1], y1 + Inches(0.29), "handoff")
    _link(slide, x[1] + sw, y1 + Inches(0.29), x[2], y1 + Inches(0.29))
    _link(slide, x[2] + sw, y1 + Inches(0.29), x[3], y2 + Inches(0.29))
    _link(slide, x[3] + sw, y2 + Inches(0.29), x[4], y1 + Inches(0.29), "sync")
    _link(slide, x[4] + sw, y1 + Inches(0.29), x[5], y1 + Inches(0.29))
    _link(slide, x[5] + sw, y1 + Inches(0.29), x[5] + Inches(1.12), y1 + Inches(0.29))

    _footer_meta(ctx, slide, "Owner: Product Manager · Output: Frozen Outline + Scope Baseline")


def slide_l1_knowledge(ctx):
    slide = _add_content_slide(ctx)
    header(slide, "L1 Process P1 — Knowledge & Template Preparation")
    _subtitle(ctx, slide, "Layer 1: 资料清洗 -> 母版选择 -> 知识注入")

    left = ctx.safe_left
    top = ctx.content_top
    width = ctx.safe_width
    lane_h = Inches(1.47)

    _lane(slide, left, top + lane_h * 0, width, lane_h, "知识运营")
    _lane(slide, left, top + lane_h * 1, width, lane_h, "设计规范")
    _lane(slide, left, top + lane_h * 2, width, lane_h, "平台服务")

    sw = Inches(1.0)
    x = _step_positions(left, width, 6, sw)
    sh = Inches(0.58)
    y0 = top + Inches(0.44)
    y1 = top + lane_h + Inches(0.44)
    y2 = top + lane_h * 2 + Inches(0.44)

    _event(slide, x[0] - Inches(0.24), y0 + Inches(0.2), start=True)
    _step(slide, x[0], y0, sw, sh, "资料采集", "KB", "<=3h")
    _step(slide, x[1], y0, sw, sh, "结构化清洗", "Parser", "<=2h")
    _step(slide, x[2], y1, sw, sh, "模板匹配", "Master Lib", "<=30m", critical=True)
    _step(slide, x[3], y1, sw, sh, "品牌规则校验", "Rule Engine", "<=20m")
    _step(slide, x[4], y2, sw, sh, "知识注入", "Vector DB", "<=25m")
    _step(slide, x[5], y2, sw, sh, "上下文发布", "Context API", "<=10m")
    _event(slide, x[5] + Inches(1.12), y2 + Inches(0.2), start=False)

    _link(slide, x[0] + sw, y0 + Inches(0.29), x[1], y0 + Inches(0.29))
    _link(slide, x[1] + sw, y0 + Inches(0.29), x[2], y1 + Inches(0.29), "handoff")
    _link(slide, x[2] + sw, y1 + Inches(0.29), x[3], y1 + Inches(0.29))
    _link(slide, x[3] + sw, y1 + Inches(0.29), x[4], y2 + Inches(0.29), "sync")
    _link(slide, x[4] + sw, y2 + Inches(0.29), x[5], y2 + Inches(0.29))
    _link(slide, x[5] + sw, y2 + Inches(0.29), x[5] + Inches(1.12), y2 + Inches(0.29))

    _footer_meta(ctx, slide, "Owner: Knowledge Ops · Output: Ready Context + Qualified Master")


def slide_l1_generation(ctx):
    slide = _add_content_slide(ctx)
    header(slide, "L1 Process P3 — Content Generation & Diagram Rendering")
    _subtitle(ctx, slide, "Layer 1: 文本生成 -> 架构制图 -> 页面编排")

    left = ctx.safe_left
    top = ctx.content_top
    width = ctx.safe_width
    lane_h = Inches(1.47)

    _lane(slide, left, top + lane_h * 0, width, lane_h, "生成引擎")
    _lane(slide, left, top + lane_h * 1, width, lane_h, "架构制图")
    _lane(slide, left, top + lane_h * 2, width, lane_h, "版式编排")

    sw = Inches(1.0)
    x = _step_positions(left, width, 6, sw)
    sh = Inches(0.58)
    y0 = top + Inches(0.44)
    y1 = top + lane_h + Inches(0.44)
    y2 = top + lane_h * 2 + Inches(0.44)

    _event(slide, x[0] - Inches(0.24), y0 + Inches(0.2), start=True)
    _step(slide, x[0], y0, sw, sh, "章节生成", "LLM", "<=4m", critical=True)
    _step(slide, x[1], y1, sw, sh, "图谱识别", "Pattern DB", "<=2m")
    _step(slide, x[2], y1, sw, sh, "架构图绘制", "Diagram Engine", "<=3m", critical=True)
    _step(slide, x[3], y2, sw, sh, "版式映射", "Layout Solver", "<=2m")
    _step(slide, x[4], y2, sw, sh, "母版合成", "PPTX Core", "<=2m")
    _step(slide, x[5], y2, sw, sh, "草稿输出", "PPTX", "<=1m")
    _event(slide, x[5] + Inches(1.12), y2 + Inches(0.2), start=False)

    _link(slide, x[0] + sw, y0 + Inches(0.29), x[1], y1 + Inches(0.29), "handoff")
    _link(slide, x[1] + sw, y1 + Inches(0.29), x[2], y1 + Inches(0.29))
    _link(slide, x[2] + sw, y1 + Inches(0.29), x[3], y2 + Inches(0.29), "sync")
    _link(slide, x[3] + sw, y2 + Inches(0.29), x[4], y2 + Inches(0.29))
    _link(slide, x[4] + sw, y2 + Inches(0.29), x[5], y2 + Inches(0.29))
    _link(slide, x[5] + sw, y2 + Inches(0.29), x[5] + Inches(1.12), y2 + Inches(0.29))

    _footer_meta(ctx, slide, "Owner: Generation Service · Output: Draft Deck + Rendered Architecture Slides")


def slide_l1_delivery(ctx):
    slide = _add_content_slide(ctx)
    header(slide, "L1 Process P4-P5 — QA, Review, Delivery & Learning")
    _subtitle(ctx, slide, "Layer 1: 质量门禁 -> 人审反馈 -> 发布归档 -> 经验回流")

    left = ctx.safe_left
    top = ctx.content_top
    width = ctx.safe_width
    lane_h = Inches(1.47)

    _lane(slide, left, top + lane_h * 0, width, lane_h, "质量保障")
    _lane(slide, left, top + lane_h * 1, width, lane_h, "评审协作")
    _lane(slide, left, top + lane_h * 2, width, lane_h, "交付运营")

    sw = Inches(1.0)
    x = _step_positions(left, width, 6, sw)
    sh = Inches(0.58)
    y0 = top + Inches(0.44)
    y1 = top + lane_h + Inches(0.44)
    y2 = top + lane_h * 2 + Inches(0.44)

    _event(slide, x[0] - Inches(0.24), y0 + Inches(0.2), start=True)
    _step(slide, x[0], y0, sw, sh, "完整性校验", "QA Engine", "<=2m", critical=True)
    _step(slide, x[1], y0, sw, sh, "风格一致性", "Style Checker", "<=2m")
    _step(slide, x[2], y1, sw, sh, "人工评审", "Review Board", "<=2h")
    _step(slide, x[3], y1, sw, sh, "反馈整合", "Workflow", "<=1h")
    _step(slide, x[4], y2, sw, sh, "正式发布", "Object Storage", "<=20m", critical=True)
    _step(slide, x[5], y2, sw, sh, "知识回流", "Learning DB", "<=30m")
    _event(slide, x[5] + Inches(1.12), y2 + Inches(0.2), start=False)

    _link(slide, x[0] + sw, y0 + Inches(0.29), x[1], y0 + Inches(0.29))
    _link(slide, x[1] + sw, y0 + Inches(0.29), x[2], y1 + Inches(0.29), "handoff")
    _link(slide, x[2] + sw, y1 + Inches(0.29), x[3], y1 + Inches(0.29))
    _link(slide, x[3] + sw, y1 + Inches(0.29), x[4], y2 + Inches(0.29), "sync")
    _link(slide, x[4] + sw, y2 + Inches(0.29), x[5], y2 + Inches(0.29))
    _link(slide, x[5] + sw, y2 + Inches(0.29), x[5] + Inches(1.12), y2 + Inches(0.29))

    _footer_meta(ctx, slide, "Owner: QA + Delivery Ops · Output: Published Deck + Feedback Signals")


def slide_l2_intake(ctx):
    slide = _add_content_slide(ctx)
    header(slide, "L2 Detail A — Intake to Scope Freeze")
    _subtitle(ctx, slide, "Detailed SOP with branch and exception (Phase P0-P2)")

    left = ctx.safe_left
    top = ctx.content_top
    width = ctx.safe_width
    lane_h = Inches(1.38)

    _lane(slide, left, top + lane_h * 0, width, lane_h, "业务方")
    _lane(slide, left, top + lane_h * 1, width, lane_h, "产品经理")
    _lane(slide, left, top + lane_h * 2, width, lane_h, "架构师")
    _lane(slide, left, top + lane_h * 3, width, lane_h, "系统")

    sw = Inches(0.92)
    x0, x1, x2, x3, x4, x5, x6 = _step_positions(left, width, 7, sw, start_pad=1.35, end_pad=1.02)
    sh = Inches(0.56)

    y_biz = top + Inches(0.41)
    y_pm = top + lane_h + Inches(0.41)
    y_arch = top + lane_h * 2 + Inches(0.41)
    y_sys = top + lane_h * 3 + Inches(0.41)

    _event(slide, x0 - Inches(0.22), y_biz + Inches(0.18), start=True)
    _step(slide, x0, y_biz, sw, sh, "提交需求", "Portal", "<=2h")
    _step(slide, x1, y_pm, sw, sh, "需求分级", "Jira", "<=1h")
    _step(slide, x2, y_pm, sw, sh, "业务澄清", "Workshop", "<=2h")
    _decision(slide, x3, y_pm + Inches(0.02), Inches(0.56), Inches(0.52), "范围清晰?")
    _step(slide, x4, y_arch, sw, sh, "约束评估", "EA Repo", "<=2h")
    _step(slide, x5, y_sys, sw, sh, "建档", "Tracker", "<=30m")
    _step(slide, x6, y_pm, sw, sh, "范围冻结", "Workflow", "<=30m", critical=True)
    _event(slide, x6 + Inches(1.02), y_pm + Inches(0.18), start=False)

    _link(slide, x0 + sw, y_biz + Inches(0.28), x1, y_pm + Inches(0.28), "handoff")
    _link(slide, x1 + sw, y_pm + Inches(0.28), x2, y_pm + Inches(0.28))
    _link(slide, x2 + sw, y_pm + Inches(0.28), x3, y_pm + Inches(0.28))
    _link(slide, x3 + Inches(0.56), y_pm + Inches(0.28), x4, y_arch + Inches(0.28), "Yes")
    _link(slide, x4 + sw, y_arch + Inches(0.28), x5, y_sys + Inches(0.28))
    _link(slide, x5 + sw, y_sys + Inches(0.28), x6, y_pm + Inches(0.28), "sync")
    _link(slide, x6 + sw, y_pm + Inches(0.28), x6 + Inches(1.02), y_pm + Inches(0.28))

    # No/Exception path
    _step(slide, x2, y_arch, sw, sh, "补充材料", "Manual", "<=4h", manual=True)
    _link(slide, x3 + Inches(0.28), y_pm + Inches(0.52), x2 + Inches(0.45), y_arch + Inches(0.02),
          label="No", dashed=True, color=C["brand_gray"])
    _link(slide, x2 + sw, y_arch + Inches(0.28), x2 + sw + Inches(0.38), y_arch + Inches(0.28),
          label="Exception", dashed=True, color=C["brand_gray"])

    _footer_meta(ctx, slide, "SLA: Intake <= 1 day · Exception path requires manual supplement")


def slide_l2_generation(ctx):
    slide = _add_content_slide(ctx)
    header(slide, "L2 Detail B — Draft Generation and Diagram Rendering")
    _subtitle(ctx, slide, "Detailed SOP with quality branch (Phase P3)")

    left = ctx.safe_left
    top = ctx.content_top
    width = ctx.safe_width
    lane_h = Inches(1.38)

    _lane(slide, left, top + lane_h * 0, width, lane_h, "生成引擎")
    _lane(slide, left, top + lane_h * 1, width, lane_h, "架构制图")
    _lane(slide, left, top + lane_h * 2, width, lane_h, "版式引擎")
    _lane(slide, left, top + lane_h * 3, width, lane_h, "质量门禁")

    sw = Inches(0.9)
    x = _step_positions(left, width, 7, sw, start_pad=1.35, end_pad=1.72)
    sh = Inches(0.56)
    y0 = top + Inches(0.41)
    y1 = top + lane_h + Inches(0.41)
    y2 = top + lane_h * 2 + Inches(0.41)
    y3 = top + lane_h * 3 + Inches(0.41)

    _event(slide, x[0] - Inches(0.22), y0 + Inches(0.18), start=True)
    _step(slide, x[0], y0, sw, sh, "章节生成", "LLM", "<=3m", critical=True)
    _step(slide, x[1], y1, sw, sh, "图模式识别", "Pattern", "<=1m")
    _step(slide, x[2], y1, sw, sh, "架构图绘制", "Diagram", "<=2m", critical=True)
    _step(slide, x[3], y2, sw, sh, "版式求解", "Layout", "<=1m")
    _step(slide, x[4], y2, sw, sh, "母版合成", "PPTX", "<=1m")
    _step(slide, x[5], y3, sw, sh, "质量预检", "QA", "<=1m")
    _decision(slide, x[6], y3 + Inches(0.02), Inches(0.52), Inches(0.50), "通过?")

    _link(slide, x[0] + sw, y0 + Inches(0.28), x[1], y1 + Inches(0.28), "handoff")
    _link(slide, x[1] + sw, y1 + Inches(0.28), x[2], y1 + Inches(0.28))
    _link(slide, x[2] + sw, y1 + Inches(0.28), x[3], y2 + Inches(0.28), "sync")
    _link(slide, x[3] + sw, y2 + Inches(0.28), x[4], y2 + Inches(0.28))
    _link(slide, x[4] + sw, y2 + Inches(0.28), x[5], y3 + Inches(0.28), "handoff")
    _link(slide, x[5] + sw, y3 + Inches(0.28), x[6], y3 + Inches(0.28))

    _step(slide, x[6] + Inches(0.72), y2, sw, sh, "输出草稿", "PPTX", "<=30s")
    _event(slide, x[6] + Inches(1.72), y2 + Inches(0.18), start=False)
    _link(slide, x[6] + Inches(0.52), y3 + Inches(0.28), x[6] + Inches(0.72), y2 + Inches(0.28), "Yes")
    _link(slide, x[6] + Inches(1.62), y2 + Inches(0.28), x[6] + Inches(1.72), y2 + Inches(0.28))

    # No path to manual repair
    _step(slide, x[4], y1, sw, sh, "人工修图", "Office", "<=30m", manual=True)
    _link(slide, x[6] + Inches(0.26), y3 + Inches(0.50), x[4] + Inches(0.45), y1 + Inches(0.02),
          label="No", dashed=True, color=C["brand_gray"])

    _footer_meta(ctx, slide, "Critical Path SLA: <= 10 min for auto-generation")


def slide_l2_delivery(ctx):
    slide = _add_content_slide(ctx)
    header(slide, "L2 Detail C — QA Review, Publish and Learning Loop")
    _subtitle(ctx, slide, "Detailed SOP with approval gate and feedback loop (Phase P4-P5)")

    left = ctx.safe_left
    top = ctx.content_top
    width = ctx.safe_width
    lane_h = Inches(1.38)

    _lane(slide, left, top + lane_h * 0, width, lane_h, "质量保障")
    _lane(slide, left, top + lane_h * 1, width, lane_h, "评审委员会")
    _lane(slide, left, top + lane_h * 2, width, lane_h, "交付运营")
    _lane(slide, left, top + lane_h * 3, width, lane_h, "学习系统")

    sw = Inches(0.9)
    x = _step_positions(left, width, 7, sw, start_pad=1.35, end_pad=1.02)
    sh = Inches(0.56)
    y0 = top + Inches(0.41)
    y1 = top + lane_h + Inches(0.41)
    y2 = top + lane_h * 2 + Inches(0.41)
    y3 = top + lane_h * 3 + Inches(0.41)

    _event(slide, x[0] - Inches(0.22), y0 + Inches(0.18), start=True)
    _step(slide, x[0], y0, sw, sh, "完整性校验", "QA", "<=1m")
    _step(slide, x[1], y0, sw, sh, "样式检查", "Style", "<=1m")
    _step(slide, x[2], y1, sw, sh, "评审会签", "Review", "<=2h", critical=True)
    _decision(slide, x[3], y1 + Inches(0.02), Inches(0.56), Inches(0.50), "批准?")
    _step(slide, x[4], y2, sw, sh, "正式发布", "Storage", "<=15m", critical=True)
    _step(slide, x[5], y2, sw, sh, "通知分发", "Notify", "<=5m")
    _step(slide, x[6], y3, sw, sh, "经验回流", "Learning DB", "<=20m")
    _event(slide, x[6] + Inches(1.02), y3 + Inches(0.18), start=False)

    _link(slide, x[0] + sw, y0 + Inches(0.28), x[1], y0 + Inches(0.28))
    _link(slide, x[1] + sw, y0 + Inches(0.28), x[2], y1 + Inches(0.28), "handoff")
    _link(slide, x[2] + sw, y1 + Inches(0.28), x[3], y1 + Inches(0.28))
    _link(slide, x[3] + Inches(0.56), y1 + Inches(0.28), x[4], y2 + Inches(0.28), "Yes")
    _link(slide, x[4] + sw, y2 + Inches(0.28), x[5], y2 + Inches(0.28))
    _link(slide, x[5] + sw, y2 + Inches(0.28), x[6], y3 + Inches(0.28), "event")
    _link(slide, x[6] + sw, y3 + Inches(0.28), x[6] + Inches(1.02), y3 + Inches(0.28))

    # rejection path
    _step(slide, x[2], y0, sw, sh, "整改回炉", "Workflow", "<=4h", manual=True)
    _link(slide, x[3] + Inches(0.28), y1 + Inches(0.50), x[2] + Inches(0.45), y0 + Inches(0.02),
          label="No", dashed=True, color=C["brand_gray"])

    # learning feedback loop to intake
    _link(slide, x[6] + Inches(0.45), y3 + Inches(0.56), x[0], y0 + Inches(0.56),
          label="feedback", dashed=True, color=C["brand_gray"])

    _footer_meta(ctx, slide, "Approval Gate + Feedback Loop closes continuous improvement cycle")


def slide_cover(ctx):
    slide = ctx.add_cover_slide()
    textbox(slide, ctx.safe_left, Inches(1.26), ctx.safe_width, Inches(1.00),
            "A-Deck", size=36, bold=True, color=C["cyan"], align=PP_ALIGN.CENTER)
    textbox(slide, ctx.safe_left, Inches(2.42), ctx.safe_width, Inches(0.70),
            "业务流程架构",
            size=26, bold=True, color=C["dark"], align=PP_ALIGN.CENTER)
    textbox(slide, ctx.safe_left + Inches(0.06), Inches(3.30), ctx.safe_width - Inches(0.12), Inches(0.80),
            "知识驱动 · 母版继承 · 端到端自动化",
            size=18, color=C["gray"], align=PP_ALIGN.CENTER)
    textbox(slide, ctx.safe_left + Inches(0.12), Inches(4.80), ctx.safe_width - Inches(0.24), Inches(0.45),
            "BA-03 Layered Business Process  |  Cloudwise Master  |  Value Chain → Swimlane SOP",
            size=13, color=C["dark"], align=PP_ALIGN.CENTER)
    textbox(slide, ctx.safe_left + Inches(0.08), Inches(6.22), ctx.safe_width - Inches(0.16), Inches(0.26),
            "From Requirement Intake to Delivery & Learning Loop",
            size=11, color=C["mid_gray"], align=PP_ALIGN.CENTER)


def slide_thanks(ctx):
    slide = ctx.add_thanks_slide()
    rrect(slide, ctx.safe_left, Inches(3.84), ctx.safe_width, Inches(1.24), C["soft_blue"], line=C["line"])
    textbox(slide, ctx.safe_left + Inches(0.2), Inches(1.45), ctx.safe_width - Inches(0.24), Inches(0.9),
            "A-Deck — 让每一份 PPT 都达到交付标准",
            size=28, bold=True, color=C["dark"], align=PP_ALIGN.CENTER)
    textbox(slide, ctx.safe_left + Inches(0.42), Inches(2.55), ctx.safe_width - Inches(0.72), Inches(0.9),
            "从需求受理到知识回流，六大阶段全流程覆盖，每个环节可追溯、可复现。",
            size=16, color=C["gray"], align=PP_ALIGN.CENTER)
    textbox(slide, ctx.safe_left + Inches(0.36), Inches(4.22), ctx.safe_width - Inches(0.62), Inches(0.5),
            "关键路径 SLA ≤ 10 分钟自动生成  |  人工评审门禁  |  经验持续回流",
            size=15, bold=True, color=C["dark"], align=PP_ALIGN.CENTER)
    stats = [
        ("Knowledge Driven", "知识库驱动全流程"),
        ("Master Based", "品牌母版保证一致性"),
        ("QA First", "质量门禁贯穿始终"),
    ]
    left = ctx.safe_left
    gap = Inches(0.18)
    width = ctx.safe_width
    box_w = (width - gap * 2) / 3
    for idx, (title, desc) in enumerate(stats):
        x = left + idx * (box_w + gap)
        rrect(slide, x, Inches(5.45), box_w, Inches(0.95), C["light"], line=C["line"])
        textbox(slide, x + Inches(0.04), Inches(5.62), box_w - Inches(0.08), Inches(0.25),
                title, size=12, bold=True, color=C["cyan"], align=PP_ALIGN.CENTER)
        textbox(slide, x + Inches(0.04), Inches(5.93), box_w - Inches(0.08), Inches(0.2),
                desc, size=10, color=C["gray"], align=PP_ALIGN.CENTER)
    footer(slide)


def my_slides(ctx):
    slide_cover(ctx)
    slide_l0(ctx)
    slide_l1_intake(ctx)
    slide_l1_knowledge(ctx)
    slide_l1_generation(ctx)
    slide_l1_delivery(ctx)
    slide_l2_intake(ctx)
    slide_thanks(ctx)


if __name__ == "__main__":
    output = build_pptx(TEMPLATE, OUTPUT, my_slides)
    verify_pptx(output, TEMPLATE)
