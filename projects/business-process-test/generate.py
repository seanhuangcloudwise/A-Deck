#!/usr/bin/env python3
"""Generate one-slide Business Process diagram (BA-03 L2 swimlane variant)."""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent / "skills" / "pptx" / "scripts"))
from pptx_lib import *  # noqa: E402,F403

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TEMPLATE = Path("/Volumes/work/Workspace/A-Deck/projects/ppt-maker-intro-cloudwise-2025"
                "/ppt-maker-agent-self-intro-cloudwise.pptx")
OUTPUT = Path(__file__).resolve().parent / "ppt-maker-business-process.pptx"


def _i(v):
    return int(v)


def _lane(slide, x, y, w, h, title):
    # lane header
    rect(slide, x, y, Inches(1.22), h, C["domain_bg"])
    textbox(slide, x + Inches(0.08), y + Inches(0.08), Inches(1.06), h - Inches(0.16),
            title, size=10, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    # lane body
    rrect(slide, x + Inches(1.24), y, w - Inches(1.24), h, C["white"], line=C["line"], adj=2500)


def _step(slide, x, y, w, h, name, system, duration, critical=False):
    edge = C["cyan"] if critical else C["line"]
    step = rrect(slide, x, y, w, h, C["cyan_2"], line=edge, adj=2600)
    step.line.width = Pt(2.0 if critical else 1.0)

    textbox(slide, x + Inches(0.07), y + Inches(0.08), w - Inches(0.14), Inches(0.26),
            name, size=9, color=C["dark"], align=PP_ALIGN.CENTER)

    # system pill
    pill = rrect(slide, x + Inches(0.18), y + Inches(0.38), w - Inches(0.36), Inches(0.15),
                 C["domain_bg"], line=C["domain_bg"], adj=2500)
    add_text(pill, system, size=7, color=C["white"], align=PP_ALIGN.CENTER)

    # duration badge
    dur = rrect(slide, x + w - Inches(0.35), y + Inches(0.02), Inches(0.30), Inches(0.13),
                C["light"], line=C["brand_gray"], adj=2000)
    add_text(dur, duration, size=6, color=C["gray"], align=PP_ALIGN.CENTER)

    if critical:
        textbox(slide, x + w - Inches(0.10), y - Inches(0.03), Inches(0.08), Inches(0.08),
                "!", size=8, bold=True, color=RGBColor(0xC0, 0x39, 0x2B), align=PP_ALIGN.CENTER)

    return step


def _decision(slide, x, y, w, h, label):
    d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, _i(x), _i(y), _i(w), _i(h))
    d.fill.solid()
    d.fill.fore_color.rgb = C["white"]
    d.line.color.rgb = C["dark"]
    d.line.width = Pt(1.2)
    textbox(slide, x + Inches(0.04), y + Inches(0.15), w - Inches(0.08), Inches(0.2),
            label, size=8, color=C["dark"], align=PP_ALIGN.CENTER)
    return d


def _event(slide, x, y, start=True):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, _i(x), _i(y), _i(Inches(0.18)), _i(Inches(0.18)))
    c.fill.solid()
    c.fill.fore_color.rgb = C["domain_bg"] if start else C["dark"]
    c.line.color.rgb = C["domain_bg"] if start else C["dark"]
    return c


def _link(slide, x1, y1, x2, y2, dashed=False, color=None, label=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _i(x1), _i(y1), _i(x2), _i(y2))
    conn.line.color.rgb = color or C["dark"]
    conn.line.width = Pt(1.4)
    if dashed:
        conn.line.dash_style = 4
    if label:
        lx = x1 + Inches(0.04)
        ly = y1 - Inches(0.14)
        textbox(slide, lx, ly, Inches(0.7), Inches(0.12), label, size=7, color=C["brand_gray"], align=PP_ALIGN.LEFT)


def my_slides(ctx):
    slide = ctx.add_content_slide()
    header(slide, "PPT Maker Agent — Business Process")
    textbox(slide, SAFE_LEFT + Inches(0.04), Inches(0.68), Inches(7.6), Inches(0.22),
            "L2 Swimlane · Owner: Product Ops · Version: v1.3", size=10, color=C["brand_gray"])

    left = SAFE_LEFT
    top = Inches(1.08)
    width = SAFE_WIDTH
    lane_h = Inches(1.38)

    _lane(slide, left, top + 0 * lane_h, width, lane_h, "业务方")
    _lane(slide, left, top + 1 * lane_h, width, lane_h, "产品/架构")
    _lane(slide, left, top + 2 * lane_h, width, lane_h, "平台/系统")
    _lane(slide, left, top + 3 * lane_h, width, lane_h, "运维/交付")

    # key x anchors
    x0 = left + Inches(1.35)
    x1 = left + Inches(2.35)
    x2 = left + Inches(3.55)
    x3 = left + Inches(4.75)
    x4 = left + Inches(5.95)
    x5 = left + Inches(7.15)
    step_w = Inches(1.0)
    step_h = Inches(0.58)

    y_biz = top + Inches(0.38)
    y_arch = top + lane_h + Inches(0.38)
    y_sys = top + lane_h * 2 + Inches(0.38)
    y_ops = top + lane_h * 3 + Inches(0.38)

    _event(slide, x0 - Inches(0.26), y_biz + Inches(0.20), start=True)

    s1 = _step(slide, x0, y_biz, step_w, step_h, "提交需求", "Portal", "<=2h")
    s2 = _step(slide, x1, y_arch, step_w, step_h, "需求澄清", "Jira", "<=4h")
    s3 = _step(slide, x2, y_arch, step_w, step_h, "架构评审", "EA Repo", "<=1d", critical=True)
    d1 = _decision(slide, x3, y_arch + Inches(0.02), Inches(0.62), Inches(0.54), "通过?")
    s4 = _step(slide, x4, y_sys, step_w, step_h, "生成草稿PPT", "PPT Maker", "<=10m", critical=True)
    s5 = _step(slide, x5, y_ops, step_w, step_h, "发布与归档", "Object Storage", "<=30m")
    _event(slide, x5 + Inches(1.12), y_ops + Inches(0.20), start=False)

    # main flow
    _link(slide, x0 + step_w, y_biz + Inches(0.29), x1, y_arch + Inches(0.29), label="handoff")
    _link(slide, x1 + step_w, y_arch + Inches(0.29), x2, y_arch + Inches(0.29))
    _link(slide, x2 + step_w, y_arch + Inches(0.29), x3, y_arch + Inches(0.29))

    # decision branches
    _link(slide, x3 + Inches(0.62), y_arch + Inches(0.29), x4, y_sys + Inches(0.29), label="Yes")
    _link(slide, x3 + Inches(0.31), y_arch + Inches(0.54), x1 + Inches(0.25), y_arch + Inches(0.58), dashed=True, color=C["brand_gray"], label="No")

    _link(slide, x4 + step_w, y_sys + Inches(0.29), x5, y_ops + Inches(0.29), label="artifact")

    # exception path: generation fail -> manual fix
    ex = _step(slide, x4, y_arch, step_w, step_h, "人工修订", "Office", "<=2h")
    _link(slide, x4 + Inches(0.1), y_sys + Inches(0.02), x4 + Inches(0.1), y_arch + step_h, dashed=True, color=C["brand_gray"], label="Exception")

    _link(slide, x5 + step_w, y_ops + Inches(0.29), x5 + Inches(1.12), y_ops + Inches(0.29))

    textbox(slide, SAFE_LEFT, Inches(7.14), SAFE_WIDTH, Inches(0.18),
            "Start: Demand Intake  |  End: Deck Published  |  Exception path shown with dashed connectors",
            size=7, color=C["brand_gray"], align=PP_ALIGN.LEFT)


if __name__ == "__main__":
    build_pptx(TEMPLATE, OUTPUT, my_slides)
    verify_pptx(OUTPUT, TEMPLATE)
