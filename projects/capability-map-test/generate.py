#!/usr/bin/env python3
"""Generate a single-slide A-Deck product capability map.

Uses Cloudwise master layout "内容" and fills title/subtitle placeholders.
Capability data is sourced from repository facts (README + product spec).
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent / "skills" / "pptx" / "scripts"))
from pptx_lib import *  # noqa: E402,F403

from pptx.util import Inches

TEMPLATE = (Path(__file__).resolve().parent.parent.parent
            / "skills" / "pptx" / "master-library"
            / "light-cloudwise-cyan" / "cloudwise-master.pptx")
OUTPUT = Path(__file__).resolve().parent / "a-deck-product-capability-map.pptx"

# A-Deck factual capability data (source: README + docs/ppt-maker-agent-spec.md)
DOMAINS = [
    {"name": "核心产品能力（README）", "caps": [
        ("学习", "GA"), ("生成", "GA"), ("质检", "GA"), ("编辑", "GA"),
    ]},
    {"name": "工作流能力（W1-W6）", "caps": [
        ("从零创建", "GA"), ("编辑现有PPT", "GA"), ("分析PPT", "GA"),
        ("模板化生成", "GA"), ("从PPT学习", "GA"), ("母版提取", "GA"),
    ]},
    {"name": "规划中扩展能力（W7-W10）", "caps": [
        ("战略路线图设计", "Planned"), ("优先级与组合平衡", "Planned"),
        ("季度执行规划", "Planned"), ("KPI回顾与再规划", "Planned"),
    ]},
]


def _add_content_slide(ctx):
    layout = layout_by_names(ctx.prs, ["内容"], 1)
    return ctx.prs.slides.add_slide(layout)


def _subtitle(slide, text):
    subtitle_ph = next(
        (s for s in slide.shapes
         if getattr(s, "is_placeholder", False) and s.placeholder_format.idx == 1),
        None,
    )
    if subtitle_ph is not None and subtitle_ph.has_text_frame:
        tf = subtitle_ph.text_frame
        tf.clear()
        tf.paragraphs[0].text = text
        return
    textbox(slide, SAFE_LEFT + Inches(0.04), Inches(0.68), Inches(8.8), Inches(0.22),
            text, size=10, color=C["brand_gray"])


def my_slides(ctx):
    slide = _add_content_slide(ctx)

    header(slide, "A-Deck Product Capability Map")
    _subtitle(slide, "Cloudwise 内容版式 · 基于仓库事实能力清单（README + W1-W10）")

    render_capability_map(slide, DOMAINS)

    textbox(slide, SAFE_LEFT, SAFE_BOTTOM + Inches(0.02), SAFE_WIDTH, Inches(0.18),
            "Source: README core capabilities + docs/ppt-maker-agent-spec.md workflow contracts",
            size=7, color=C["brand_gray"])


if __name__ == "__main__":
    build_pptx(TEMPLATE, OUTPUT, my_slides)
    verify_pptx(OUTPUT, TEMPLATE)
