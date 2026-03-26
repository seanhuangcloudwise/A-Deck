#!/usr/bin/env python3
"""
GTM A-Deck 完整演示PPT生成器

架构：
  generate.py (编排器)
    ↓ 读取配置
    ↓ 加载数据
    ↓ 依次调用装载器
    ↓
  loaders/ (渲染引擎)
    gm_01, gm_02, ..., gm_35 (35个装载器，每个接收data)
    ↓ 数据驱动渲染
    ↓
  shared/ (共用工具)
    renderer_utils.py (通用render函数)

工作流：
  1. 从 data/config_template.yaml 读取所有业务数据
  2. 遍历装载器列表，为每个装载器传入对应的data
  3. 每个装载器完全参数化，无硬编码数据
"""

import sys
from pathlib import Path
import yaml
import importlib

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent / "skills" / "pptx" / "scripts"))
from pptx_lib import *  # noqa: E402,F403
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent.parent
TEMPLATE = ROOT / "skills" / "pptx" / "master-library" / "light-cloudwise-purple" / "cloudwise-master.pptx"
TEMPLATE_SPEC = ROOT / "skills" / "pptx" / "master-library" / "light-cloudwise-purple" / "cloudwise-spec.yaml"
OUTPUT = Path(__file__).resolve().parent / "a-deck-gtm-full-demo.pptx"
CONFIG_FILE = Path(__file__).resolve().parent / "data" / "config_template.yaml"

# 技能库位置
SKILLS_ROOT = ROOT / "skills"
LOADERS_DIR = SKILLS_ROOT / "pptx" / "gtm-architecture" / "loaders"
SHARED_DIR = SKILLS_ROOT / "pptx" / "shared"

sys.path.insert(0, str(LOADERS_DIR.parent))
sys.path.insert(0, str(SHARED_DIR))
from renderer_utils import configure_theme

# ============================================================================
# 数据加载
# ============================================================================

def load_config(config_path):
    """加载YAML配置文件"""
    with open(config_path, 'r', encoding='utf-8') as f:
        return yaml.safe_load(f)


def add_content_slide(ctx):
    """添加内容幻灯片"""
    layout = layout_by_names(ctx.prs, ["标题和内容", "内容", "內容"], 2)
    return ctx.prs.slides.add_slide(layout)


# ============================================================================
# 装载器加载
# ============================================================================

LOADER_REGISTRY = {}


def register_loader(gm_name, load_function):
    """注册装载器"""
    LOADER_REGISTRY[gm_name] = load_function


def load_all_loaders():
    """动态加载所有装载器模块（从skills库中加载）"""
    # 确保shared在sys.path中
    if str(SHARED_DIR) not in sys.path:
        sys.path.insert(0, str(SHARED_DIR))
    if str(LOADERS_DIR.parent) not in sys.path:
        sys.path.insert(0, str(LOADERS_DIR.parent))
    
    # 找到所有 gm_*.py 文件
    loader_files = sorted(LOADERS_DIR.glob("gm_*.py"))
    
    print(f"[Info] 找到 {len(loader_files)} 个装载器模块")
    
    for loader_file in loader_files:
        module_name = loader_file.stem  # 例: gm_01_positioning_statement
        
        try:
            # 动态导入模块
            spec = importlib.util.spec_from_file_location(module_name, loader_file)
            module = importlib.util.module_from_spec(spec)
            
            # 在执行前确保模块的sys.path正确
            if str(SHARED_DIR) not in module.__dict__.get('__path__', []):
                module.__path__ = [str(SHARED_DIR)] + list(module.__dict__.get('__path__', []))
            
            spec.loader.exec_module(module)
            
            # 查找 load_slide 函数
            if hasattr(module, 'load_slide'):
                register_loader(module_name, module.load_slide)
                print(f"  ✓ {module_name}")
            else:
                print(f"  ✗ {module_name} (缺少 load_slide 函数)")
        except Exception as e:
            print(f"  ✗ {module_name} (加载失败: {e})")


# ============================================================================
# PPT生成 - 数据驱动的编排
# ============================================================================

def my_slides(ctx):
    """
    主幻灯片生成函数
    
    结构：
    1. 封面 (1页)
    2. G1 价值主张 → GM-01~03
    3. G2 市场与竞争 → GM-04~06
    4. G3 解决方案架构 → GM-07~09
    5. G4 证明与ROI → GM-10~12
    6. G5 GTM战略 → GM-13~15
    7. G6 品类与分析师 → GM-16~17
    8. G7 产品特性表达 → GM-18~25
    9. G8 价值实现 → GM-26~35
    10. 封底 (1页)
    """
    
    configure_theme(ctx.colors, ctx.template_spec)
    data = load_config(CONFIG_FILE)
    
    # 1. 封面
    add_cover_slide(ctx)
    
    # 2-9. 各分组章节
    groups = [
        ("G1 价值主张", "g1_value_proposition", [
            "gm_01_positioning_statement",
            "gm_02_value_pyramid",
            "gm_03_before_after_comparison",
        ]),
        ("G2 市场与竞争", "g2_market_competition", [
            "gm_04_market_positioning_matrix",
            "gm_05_competitive_comparison_matrix",
            "gm_06_market_ecosystem_map",
        ]),
        ("G3 解决方案架构", "g3_solution_architecture", [
            "gm_07_solution_reference_architecture",
            "gm_08_customer_journey_touchpoints",
            "gm_09_use_case_scenario",
        ]),
        ("G4 证明与ROI", "g4_proof_roi", [
            "gm_10_roi_business_case",
            "gm_11_kpi_dashboard_mockup",
            "gm_12_customer_success_metrics",
        ]),
        ("G5 GTM战略", "g5_gtm_strategy", [
            "gm_13_gtm_motion_diagram",
            "gm_14_icp_segmentation_map",
            "gm_15_sales_playbook_flow",
        ]),
        ("G6 品类与分析师", "g6_category_analyst", [
            "gm_16_category_creation",
            "gm_17_analyst_briefing_framework",
        ]),
        ("G7 产品特性表达", "g7_product_feature", [
            "gm_18_feature_capability_matrix",
            "gm_19_feature_differentiation_radar",
            "gm_20_feature_usecase_mapping",
            "gm_21_feature_depth_ladder",
            "gm_22_unique_mechanism_diagram",
            "gm_23_feature_adoption_funnel",
            "gm_24_feature_release_timeline",
            "gm_25_feature_proof_card",
        ]),
        ("G8 价值实现", "g8_value_realization", [
            "gm_26_value_driver_tree",
            "gm_27_capability_outcome_trace_matrix",
            "gm_28_time_to_value_curve",
            "gm_29_baseline_target_kpi_table",
            "gm_30_cost_of_inaction_table",
            "gm_31_benefit_realization_roadmap",
            "gm_32_persona_value_map",
            "gm_33_risk_reduction_heatmap",
            "gm_34_proof_evidence_ladder",
            "gm_35_assumption_sensitivity_table",
        ]),
    ]
    
    for group_title, group_key, loader_names in groups:
        # 添加章节隔断页
        add_chapter_slide(ctx, group_title)
        
        # 为该分组的每个图添加内容页
        for loader_name in loader_names:
            if loader_name in LOADER_REGISTRY and group_key in data:
                # 直接使用loader_name作为key在YAML中查找
                if loader_name in data[group_key]:
                    loader_func = LOADER_REGISTRY[loader_name]
                    slide_data = data[group_key][loader_name]
                    
                    print(f"[Slide] 正在生成 {loader_name}...")
                    try:
                        loader_func(ctx, slide_data)
                        print(f"  ✓ 成功")
                    except Exception as e:
                        print(f"  ✗ 生成失败: {e}")
                        import traceback
                        traceback.print_exc()
                else:
                    print(f"[Warn] {loader_name} 在配置中未找到")
            else:
                if loader_name not in LOADER_REGISTRY:
                    print(f"[Warn] {loader_name} 未在装载器中注册")
                if group_key not in data:
                    print(f"[Warn] {group_key} 在配置中未找到")
    
    # 9. 封底
    add_back_slide(ctx)


def add_cover_slide(ctx):
    """添加封面"""
    layout = layout_by_names(ctx.prs, ["标题幻灯片", "封面"], 1)
    slide = ctx.prs.slides.add_slide(layout)
    
    title_ph = next(
        (s for s in slide.shapes
         if getattr(s, "is_placeholder", False) and s.placeholder_format.idx == 0),
        None,
    )
    if title_ph and title_ph.has_text_frame:
        tf = title_ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "GTM Complete Diagram Demo"
        p.font.size = Pt(54)
        p.font.bold = True
    
    # 副标题
    subtitle_ph = next(
        (s for s in slide.shapes
         if getattr(s, "is_placeholder", False) and s.placeholder_format.idx == 1),
        None,
    )
    if subtitle_ph and subtitle_ph.has_text_frame:
        tf = subtitle_ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "A-Deck Platform - 企业PPT生成全景"


def add_chapter_slide(ctx, group_title):
    """添加章节隔断页"""
    layout = layout_by_names(ctx.prs, ["节标题", "刻度线", "節標題"], 0)
    
    slide = ctx.prs.slides.add_slide(layout)
    
    title_ph = next(
        (s for s in slide.shapes
         if getattr(s, "is_placeholder", False) and s.placeholder_format.idx == 0),
        None,
    )
    if title_ph and title_ph.has_text_frame:
        tf = title_ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = group_title
        p.font.size = Pt(44)
        p.font.bold = True


def add_back_slide(ctx):
    """添加封底"""
    layout = layout_by_names(ctx.prs, ["末尾幻灯片", "谢", "感谢"], 10)
    slide = ctx.prs.slides.add_slide(layout)
    
    title_ph = next(
        (s for s in slide.shapes
         if getattr(s, "is_placeholder", False) and s.placeholder_format.idx == 0),
        None,
    )
    if title_ph and title_ph.has_text_frame:
        tf = title_ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.font.size = Pt(54)
        p.font.bold = True


# ============================================================================
# 主执行
# ============================================================================

if __name__ == "__main__":
    print("=" * 60)
    print("GTM A-Deck 完整演示PPT 生成器")
    print("=" * 60)
    
    # 1. 加载所有装载器
    print("\n[Step 1] 加载装载器...")
    load_all_loaders()
    print(f"已注册 {len(LOADER_REGISTRY)} 个装载器\n")
    
    # 2. 生成PPT
    print("[Step 2] 生成PPT...")
    try:
        build_pptx(TEMPLATE, OUTPUT, my_slides, TEMPLATE_SPEC)
        print(f"✓ PPT已生成: {OUTPUT}\n")
        
        # 3. 验证
        print("[Step 3] 验证...")
        verify_pptx(OUTPUT, TEMPLATE)
        print("✓ 验证完成")
        
    except Exception as e:
        print(f"✗ 生成失败: {e}")
        import traceback
        traceback.print_exc()
