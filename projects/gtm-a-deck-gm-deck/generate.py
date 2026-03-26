#!/usr/bin/env python3
"""Generate A-Deck GTM deck (GM-01/02/05/06/08) on Cloudwise master layout."""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent / "skills" / "pptx" / "scripts"))
from pptx_lib import *  # noqa: E402,F403

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent.parent
TEMPLATE = ROOT / "skills" / "pptx" / "master-library" / "light-cloudwise-cyan" / "cloudwise-master.pptx"
OUTPUT = Path(__file__).resolve().parent / "a-deck-gtm-gm01-02-05-06-08.pptx"


def add_content_slide(ctx):
    layout = layout_by_names(ctx.prs, ["内容", "內容"], 1)
    return ctx.prs.slides.add_slide(layout)


def set_subtitle(slide, text):
    subtitle_ph = next(
        (s for s in slide.shapes
         if getattr(s, "is_placeholder", False) and s.placeholder_format.idx == 1),
        None,
    )
    if subtitle_ph is not None and subtitle_ph.has_text_frame:
        tf = subtitle_ph.text_frame
        tf.clear()
        tf.paragraphs[0].text = text
    else:
        textbox(slide, SAFE_LEFT + Inches(0.04), Inches(0.70), Inches(8.6), Inches(0.25),
                text, size=10, color=C["brand_gray"])


def gm01_positioning(ctx):
    slide = add_content_slide(ctx)
    header(slide, "GM-01 Positioning Statement | A-Deck 市场定位")
    set_subtitle(slide, "For 产品/营销团队 | 结构化定位链路：Target -> Problem -> Product -> Differentiator -> Category")

    items = [
        ("Target", "产品团队与营销团队\n(200-10K人规模企业)", C["cyan_3"], C["dark"]),
        ("Problem", "PPT制作慢且风格不一致\n手工协作成本高", C["light"], C["dark"]),
        ("Product", "A-Deck\n企业级PPT生成平台", C["cyan"], C["white"]),
        ("Differentiator", "模板智能 + AI生成\n品牌规范自动约束", C["domain_bg"], C["white"]),
        ("Category", "Enterprise PPT\nAutomation", C["white"], C["cyan"]),
    ]

    y = Inches(2.0)
    h = Inches(1.05)
    left = SAFE_LEFT + Inches(0.02)
    gap = Inches(0.12)
    box_w = (SAFE_WIDTH - Inches(0.10) - gap * 4) / 5

    for i, (label, content, fill, text_color) in enumerate(items):
        x = left + i * (box_w + gap)
        line = C["cyan"] if label == "Category" else C["line"]
        rrect(slide, x, y, box_w, h, fill, line=line)
        textbox(slide, x, y - Inches(0.19), box_w, Inches(0.16), label,
                size=8, color=C["brand_gray"], align=PP_ALIGN.CENTER)
        textbox(slide, x + Inches(0.04), y + Inches(0.22), box_w - Inches(0.08), Inches(0.62),
                content, size=10, bold=(label == "Product"), color=text_color,
                align=PP_ALIGN.CENTER)
        if i < 4:
            textbox(slide, x + box_w + Inches(0.01), y + Inches(0.43), Inches(0.09), Inches(0.2),
                    ">", size=13, bold=True, color=C["cyan"], align=PP_ALIGN.CENTER)

    rrect(slide, SAFE_LEFT + Inches(0.02), Inches(3.45), SAFE_WIDTH - Inches(0.10), Inches(1.25), C["soft_blue"], line=C["line"])
    textbox(
        slide,
        SAFE_LEFT + Inches(0.18), Inches(3.62), SAFE_WIDTH - Inches(0.42), Inches(0.85),
        "For product and marketing teams that need fast, brand-compliant enterprise decks,\n"
        "A-Deck is an Enterprise PPT Automation platform that generates delivery-ready slides.\n"
        "Unlike generic design tools or manual workflows, A-Deck combines template intelligence, AI generation, and governance.",
        size=11, color=C["dark"],
    )


def gm02_value_pyramid(ctx):
    slide = add_content_slide(ctx)
    header(slide, "GM-02 Value Pyramid | A-Deck 价值金字塔")
    set_subtitle(slide, "功能价值 -> 业务价值 -> 战略价值（CTO / CFO / CEO）")

    center = SAFE_LEFT + SAFE_WIDTH / 2
    y0 = Inches(4.95)

    bottom_w = Inches(7.8)
    mid_w = Inches(5.7)
    top_w = Inches(3.7)

    rrect(slide, center - bottom_w / 2, y0 - Inches(1.45), bottom_w, Inches(1.45), C["cyan_2"], line=C["line"])
    add_text(slide.shapes[-1], "Functional Value | 5分钟初稿生成 / 品牌规范自动套用 / 多版式输出", 12, bold=True, color=C["dark"])

    rrect(slide, center - mid_w / 2, y0 - Inches(2.55), mid_w, Inches(1.10), C["cyan"], line=C["line"])
    add_text(slide.shapes[-1], "Business Value | 交付周期 Days->Hours / 返工减少 / 质量一致", 12, bold=True, color=C["white"])

    rrect(slide, center - top_w / 2, y0 - Inches(3.35), top_w, Inches(0.80), C["domain_bg"], line=C["line"])
    add_text(slide.shapes[-1], "Strategic Value | 更快GTM与规模化内容运营", 11, bold=True, color=C["white"])

    rrect(slide, SAFE_LEFT + Inches(0.04), Inches(1.55), Inches(2.0), Inches(0.56), C["light"], line=C["line"])
    textbox(slide, SAFE_LEFT + Inches(0.12), Inches(1.70), Inches(1.85), Inches(0.24), "CTO关注: 生成效率", size=9, color=C["dark"])

    rrect(slide, SAFE_LEFT + Inches(4.10), Inches(1.55), Inches(2.2), Inches(0.56), C["soft_blue"], line=C["line"])
    textbox(slide, SAFE_LEFT + Inches(4.18), Inches(1.70), Inches(2.0), Inches(0.24), "CFO关注: 成本与复用", size=9, color=C["dark"])

    rrect(slide, SAFE_LEFT + Inches(7.0), Inches(1.55), Inches(2.2), Inches(0.56), C["soft_yellow"], line=C["line"])
    textbox(slide, SAFE_LEFT + Inches(7.08), Inches(1.70), Inches(2.0), Inches(0.24), "CEO关注: 市场响应速度", size=9, color=C["dark"])

    textbox(slide, SAFE_LEFT, SAFE_BOTTOM + Inches(0.01), SAFE_WIDTH, Inches(0.18),
            "示例指标：准备周期缩短70%，统一模板覆盖率95%（演示样例）", size=7, color=C["brand_gray"])


def draw_matrix_table(slide, title, subtitle, headers, rows, y_start=Inches(1.45)):
    header(slide, title)
    set_subtitle(slide, subtitle)

    left = SAFE_LEFT + Inches(0.02)
    total_w = SAFE_WIDTH - Inches(0.08)
    criteria_w = Inches(2.8)
    col_w = (total_w - criteria_w) / (len(headers) - 1)
    row_h = Inches(0.42)

    x = left
    for c, h in enumerate(headers):
        w = criteria_w if c == 0 else col_w
        fill = C["cyan"] if h == "A-Deck" else C["domain_bg"]
        rect(slide, x, y_start, w, row_h, fill)
        textbox(slide, x + Inches(0.03), y_start + Inches(0.08), w - Inches(0.06), Inches(0.22),
                h, size=9, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        x += w

    for r, row in enumerate(rows):
        y = y_start + row_h * (r + 1)
        x = left
        is_group = row[0].startswith("[")
        for c, cell in enumerate(row):
            w = criteria_w if c == 0 else col_w
            if is_group:
                rect(slide, x, y, w, row_h, C["domain_bg"])
                txt_color = C["white"]
                txt = cell.strip("[]") if c == 0 else ""
                bold = True
            else:
                row_fill = C["soft_blue"] if (r % 2 == 0 and c == 0) else C["white"]
                rect(slide, x, y, w, row_h, row_fill, line=C["line"])
                txt_color = C["dark"]
                txt = cell
                bold = c == 0
            textbox(slide, x + Inches(0.03), y + Inches(0.07), w - Inches(0.06), Inches(0.22),
                    txt, size=8.5, bold=bold, color=txt_color,
                    align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)
            x += w


def gm05_competition_1(ctx):
    slide = add_content_slide(ctx)
    headers = ["Capability", "A-Deck", "手工设计", "通用AI写作", "基础模板工具"]
    rows = [
        ["[基础能力]", "", "", "", ""],
        ["品牌母版继承", "✓", "~", "✗", "~"],
        ["结构化提纲先行", "✓", "✗", "~", "✗"],
        ["自动图形匹配", "✓", "✗", "✗", "~"],
        ["多页一致编排", "✓", "~", "✗", "~"],
        ["[生成效率]", "", "", "", ""],
        ["5-10分钟首版", "✓", "✗", "~", "✗"],
        ["迭代改稿速度", "✓", "✗", "~", "~"],
    ]
    draw_matrix_table(
        slide,
        "GM-05 Competitive Comparison (I) | 基础与效率",
        "评估时间: 2026-03 | ✓ Full / ~ Partial / ✗ Not available",
        headers,
        rows,
    )


def gm05_competition_2(ctx):
    slide = add_content_slide(ctx)
    headers = ["Enterprise Criteria", "A-Deck", "手工设计", "通用AI写作", "基础模板工具"]
    rows = [
        ["[企业治理]", "", "", "", ""],
        ["品牌一致性约束", "✓", "✗", "✗", "~"],
        ["可复用知识沉淀", "✓", "✗", "~", "✗"],
        ["项目级模板治理", "✓", "✗", "✗", "~"],
        ["[业务交付]", "", "", "", ""],
        ["方案型PPT适配", "✓", "~", "✗", "~"],
        ["技术/产品叙事能力", "✓", "~", "~", "✗"],
        ["交付稳定性", "✓", "~", "✗", "~"],
    ]
    draw_matrix_table(
        slide,
        "GM-05 Competitive Comparison (II) | 企业治理与交付",
        "方法: 公开能力比对 + 场景化评估（非第三方审计）",
        headers,
        rows,
    )


def gm06_ecosystem(ctx):
    slide = add_content_slide(ctx)
    header(slide, "GM-06 Market Ecosystem Map | A-Deck 生态位")
    set_subtitle(slide, "Category Grid: Content Platforms / Template & Design / AI Tools / Enterprise Governance")

    left = SAFE_LEFT + Inches(0.02)
    top = Inches(1.55)
    col_gap = Inches(0.1)
    cols = 4
    col_w = (SAFE_WIDTH - Inches(0.08) - col_gap * (cols - 1)) / cols
    box_h = Inches(4.8)

    groups = [
        ("Content Platforms", ["Confluence", "Notion", "SharePoint"]),
        ("Template & Design", ["Canva", "Figma", "PowerPoint"]),
        ("AI Tools", ["ChatGPT", "Copilot", "Jasper"]),
        ("Enterprise Governance", ["Workato", "Workiva", "SSO/Policy"]),
    ]

    for i, (name, vendors) in enumerate(groups):
        x = left + i * (col_w + col_gap)
        rrect(slide, x, top, col_w, box_h, C["light"], line=C["line"])
        rect(slide, x, top, col_w, Inches(0.45), C["domain_bg"])
        textbox(slide, x + Inches(0.05), top + Inches(0.10), col_w - Inches(0.10), Inches(0.22),
                name, size=9, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        for j, vendor in enumerate(vendors):
            yy = top + Inches(0.62) + j * Inches(0.42)
            rrect(slide, x + Inches(0.08), yy, col_w - Inches(0.16), Inches(0.30), C["white"], line=C["line"])
            textbox(slide, x + Inches(0.12), yy + Inches(0.06), col_w - Inches(0.24), Inches(0.18),
                    vendor, size=8.5, color=C["dark"], align=PP_ALIGN.CENTER)

    center_x = left + col_w + col_gap + (col_w + col_gap) / 2
    center_w = Inches(2.8)
    center_y = Inches(4.05)
    rrect(slide, center_x, center_y, center_w, Inches(1.05), C["cyan"], line=C["cyan"])
    textbox(slide, center_x + Inches(0.05), center_y + Inches(0.24), center_w - Inches(0.1), Inches(0.5),
            "A-Deck\nTemplate Intelligence + AI Generation + Governance", size=10,
            bold=True, color=C["white"], align=PP_ALIGN.CENTER)

    textbox(slide, SAFE_LEFT, SAFE_BOTTOM + Inches(0.01), SAFE_WIDTH, Inches(0.18),
            "Landscape as of 2026-03 | Inclusion: mainstream enterprise content/design/AI/governance tools", size=7, color=C["brand_gray"])


def journey_rows(slide, stages, actions, touchpoints, y, row_h):
    left = SAFE_LEFT + Inches(0.02)
    total_w = SAFE_WIDTH - Inches(0.08)
    col_gap = Inches(0.08)
    col_w = (total_w - col_gap * (len(stages) - 1)) / len(stages)

    for i, stage in enumerate(stages):
        x = left + i * (col_w + col_gap)
        rect(slide, x, y, col_w, Inches(0.38), C["domain_bg"])
        textbox(slide, x + Inches(0.02), y + Inches(0.08), col_w - Inches(0.04), Inches(0.22),
                stage, size=8.5, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

        rrect(slide, x, y + Inches(0.43), col_w, row_h, C["white"], line=C["line"])
        textbox(slide, x + Inches(0.04), y + Inches(0.50), col_w - Inches(0.08), row_h - Inches(0.1),
                actions[i], size=8.5, color=C["dark"])

        rrect(slide, x, y + Inches(0.43) + row_h + Inches(0.07), col_w, row_h, C["cyan_3"], line=C["line"])
        textbox(slide, x + Inches(0.04), y + Inches(0.50) + row_h + Inches(0.07), col_w - Inches(0.08), row_h - Inches(0.1),
                touchpoints[i], size=8.5, color=C["dark"])


def gm08_journey_part1(ctx):
    slide = add_content_slide(ctx)
    header(slide, "GM-08 Customer Journey (I) | 需求到首版生成")
    set_subtitle(slide, "角色: 产品经理 / 市场经理 | 行程: Briefing -> Planning -> Generation")

    stages = ["Briefing", "Planning", "Generation"]
    actions = [
        "需求出现:\n下周需对外方案PPT\n痛点: 时间紧、素材散",
        "梳理结构:\n统一叙事主线与章节\n痛点: 模板选择困难",
        "生成首版:\n快速形成可评审草稿\n结果: 准备周期显著缩短",
    ]
    touchpoints = [
        "A-Deck触点:\n输入主题/受众/目标",
        "A-Deck触点:\n提纲建议 + 版式匹配",
        "A-Deck触点:\n按母版自动编排生成",
    ]
    journey_rows(slide, stages, actions, touchpoints, Inches(1.55), Inches(1.5))


def gm08_journey_part2(ctx):
    slide = add_content_slide(ctx)
    header(slide, "GM-08 Customer Journey (II) | 评审迭代到发布")
    set_subtitle(slide, "角色: 市场/销售/管理层 | 行程: Review -> Revise -> Publish")

    stages = ["Review", "Revise", "Publish"]
    actions = [
        "多角色评审:\n聚焦信息准确性\n痛点: 版本分叉",
        "集中修改:\n消除风格与结构偏差\n痛点: 返工频繁",
        "发布与复用:\n用于客户沟通/内部汇报\n结果: 交付可复用",
    ]
    touchpoints = [
        "A-Deck触点:\n审阅标注与意见汇总",
        "A-Deck触点:\n一致性修订与版本收敛",
        "A-Deck触点:\n导出交付 + 知识沉淀",
    ]
    journey_rows(slide, stages, actions, touchpoints, Inches(1.55), Inches(1.5))


def my_slides(ctx):
    gm01_positioning(ctx)
    gm02_value_pyramid(ctx)
    gm05_competition_1(ctx)
    gm05_competition_2(ctx)
    gm06_ecosystem(ctx)
    gm08_journey_part1(ctx)
    gm08_journey_part2(ctx)


if __name__ == "__main__":
    out = build_pptx(TEMPLATE, OUTPUT, my_slides)
    verify_pptx(out, TEMPLATE)
    print(f"Generated: {out}")
