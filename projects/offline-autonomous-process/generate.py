#!/usr/bin/env python3
"""
断网自主作业流程图
场景：下发任务后给机器人断网，可自主作业不中断
  - 本地能处理 → 本地处理
  - 本地不能处理 → 判断素材合格性 → 存于本地，等待网络恢复

母版：light-cloudwise-purple
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent / "skills" / "pptx" / "scripts"))
from pptx_lib import *  # noqa: E402,F403

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TEMPLATE = (Path(__file__).resolve().parent.parent.parent
            / "skills" / "pptx" / "master-library"
            / "light-cloudwise-purple" / "cloudwise-master.pptx")
OUTPUT = Path(__file__).resolve().parent / "offline-autonomous-process.pptx"


def _i(v):
    return int(v)


def _rect(slide, x, y, w, h, fill, line=None, radius=2400):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _i(x), _i(y), _i(w), _i(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1.2)
    else:
        s.line.fill.background()
    adj = s.adjustments
    if len(adj) > 0:
        adj[0] = radius / 100000
    return s


def _diamond(slide, x, y, w, h, fill, line):
    s = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, _i(x), _i(y), _i(w), _i(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(1.4)
    return s


def _oval(slide, x, y, d, fill, line):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, _i(x), _i(y), _i(d), _i(d))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(1.2)
    return s


def _conn(slide, x1, y1, x2, y2, color, dashed=False, width_pt=1.3):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _i(x1), _i(y1), _i(x2), _i(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width_pt)
    if dashed:
        c.line.dash_style = 4
    return c


def _label(slide, x, y, w, h, text, size=8, color=None, bold=False, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    from pptx.util import Pt as Pt2
    run = p.add_run()
    run.text = text
    run.font.size = Pt2(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tb


def _arrowhead(slide, x, y, direction="down", color=None, size=Inches(0.07)):
    """Draw a small triangle arrowhead."""
    from pptx.enum.shapes import MSO_SHAPE
    shape_map = {
        "down":  MSO_SHAPE.ISOSCELES_TRIANGLE,
        "right": MSO_SHAPE.ISOSCELES_TRIANGLE,
    }
    s = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, _i(x - size / 2), _i(y), _i(size), _i(size))
    s.fill.solid()
    s.fill.fore_color.rgb = color or RGBColor(0x44, 0x44, 0x44)
    s.line.fill.background()
    if direction == "right":
        s.rotation = 90
    return s


def build(ctx):
    slide = ctx.add_content_slide()
    C = ctx.colors

    header(slide, "断网自主作业流程", "任务下发后断网 · 机器人本地自主决策与处理", 1)

    # ── 布局参数 ──────────────────────────────────────────────
    SL = ctx.safe_left
    ST = int(ctx.content_top + Inches(0.3))
    SW = ctx.safe_width
    SB = int(ctx.safe_bottom)

    # 总可用高度
    total_h = SB - ST - int(Inches(0.1))

    # 节点尺寸
    NW = int(Inches(1.6))   # 普通节点宽
    NH = int(Inches(0.52))  # 普通节点高
    DW = int(Inches(1.5))   # 菱形宽
    DH = int(Inches(0.58))  # 菱形高
    EV = int(Inches(0.22))  # 事件圆直径

    def _rgb(c):
        """Ensure value is RGBColor — handles tuple (r,g,b), int, or RGBColor."""
        if isinstance(c, RGBColor):
            return c
        if isinstance(c, (tuple, list)):
            return RGBColor(int(c[0]), int(c[1]), int(c[2]))
        if isinstance(c, int):
            return RGBColor((c >> 16) & 0xFF, (c >> 8) & 0xFF, c & 0xFF)
        return c

    # 颜色
    col_primary   = _rgb(C.get("primary",    RGBColor(0x5B, 0x2D, 0x8E)))
    col_secondary = _rgb(C.get("secondary",  RGBColor(0x7C, 0x3E, 0xB0)))
    col_light     = _rgb(C.get("light",      RGBColor(0xF3, 0xEE, 0xFA)))
    col_dark      = _rgb(C.get("dark",       RGBColor(0x1A, 0x0A, 0x2E)))
    col_white     = RGBColor(0xFF, 0xFF, 0xFF)
    col_line      = _rgb(C.get("line",       RGBColor(0xBB, 0xBB, 0xBB)))
    col_warn      = RGBColor(0xE6, 0x7E, 0x22)   # 橙色：离线/断网场景
    col_ok        = RGBColor(0x27, 0xAE, 0x60)   # 绿色：处理成功
    col_store     = RGBColor(0x29, 0x80, 0xB9)   # 蓝色：存储

    # ── 泳道背景 ──────────────────────────────────────────────
    # 3条横向泳道：触发 | 决策 | 结果
    lane_h = total_h // 3
    lane_labels = ["触发 & 断网", "本地决策", "处理 & 存储"]
    lane_colors = [
        RGBColor(0xFA, 0xF0, 0xFF),
        RGBColor(0xF5, 0xF0, 0xFF),
        RGBColor(0xEE, 0xEE, 0xFF),
    ]

    for i, (lbl, lc) in enumerate(zip(lane_labels, lane_colors)):
        ly = ST + i * lane_h
        # 泳道背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    _i(SL), _i(ly), _i(SW), _i(lane_h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = lc
        bg.line.color.rgb = col_line
        bg.line.width = Pt(0.5)
        # 泳道标签竖条
        label_w = int(Inches(0.9))
        lr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     _i(SL), _i(ly), _i(label_w), _i(lane_h))
        lr.fill.solid()
        lr.fill.fore_color.rgb = col_primary
        lr.line.fill.background()
        _label(slide, SL + Inches(0.05), ly + lane_h // 2 - int(Inches(0.15)),
               Inches(0.8), Inches(0.3), lbl, size=8, bold=True, color=col_white)

    # ── 内容区 X 起点（泳道标签右侧）──────────────────────────
    cx0 = int(SL + Inches(1.0))
    cw  = int(SW - Inches(1.0))

    # ── Lane 0: 触发 & 断网 ────────────────────────────────────
    L0_cy = ST + lane_h // 2  # 泳道中心Y

    # Step 0: 开始事件
    s0x = cx0 + int(Inches(0.2))
    s0y = L0_cy - EV // 2
    _oval(slide, s0x, s0y, EV, col_primary, col_primary)

    # Step 1: 下发任务
    s1x = s0x + EV + int(Inches(0.22))
    s1y = L0_cy - NH // 2
    s1 = _rect(slide, s1x, s1y, NW, NH, col_primary, col_primary)
    add_text(s1, "下发任务", size=9, bold=True, color=col_white, align=PP_ALIGN.CENTER)

    # Step 2: 断网（橙色警示）
    s2x = s1x + NW + int(Inches(0.22))
    s2y = L0_cy - NH // 2
    s2 = _rect(slide, s2x, s2y, NW, NH, col_warn, col_warn)
    add_text(s2, "断网\n（拔网线）", size=9, bold=True, color=col_white, align=PP_ALIGN.CENTER)

    # Step 3: 机器人接收任务
    s3x = s2x + NW + int(Inches(0.22))
    s3y = L0_cy - NH // 2
    s3 = _rect(slide, s3x, s3y, NW, NH, col_secondary, col_secondary)
    add_text(s3, "机器人\n接收任务", size=9, bold=True, color=col_white, align=PP_ALIGN.CENTER)

    # 连线 Lane 0
    _conn(slide, s0x + EV, L0_cy, s1x, L0_cy, col_dark)
    _conn(slide, s1x + NW, L0_cy, s2x, L0_cy, col_dark)
    _conn(slide, s2x + NW, L0_cy, s3x, L0_cy, col_dark)

    # 从 Step3 向下到 Lane 1
    s3_cx = s3x + NW // 2
    _conn(slide, s3_cx, s3y + NH, s3_cx, ST + lane_h, col_dark)

    # ── Lane 1: 本地决策 ───────────────────────────────────────
    L1_top = ST + lane_h
    L1_cy  = L1_top + lane_h // 2

    # 决策菱形：本地能处理？
    d1x = s3_cx - DW // 2
    d1y = L1_cy - DH // 2
    d1 = _diamond(slide, d1x, d1y, DW, DH, col_white, col_dark)
    _label(slide, d1x + int(Inches(0.12)), d1y + int(Inches(0.17)),
           DW - int(Inches(0.24)), int(Inches(0.26)),
           "本地能\n处理？", size=8, bold=True, color=col_dark)

    # 是 → 向右到 Lane 2（本地处理分支）
    yes_x = d1x + DW
    yes_cy = L1_cy
    _conn(slide, yes_x, yes_cy, yes_x + int(Inches(0.8)), yes_cy, col_ok, width_pt=1.5)
    _label(slide, yes_x + int(Inches(0.04)), yes_cy - int(Inches(0.18)),
           Inches(0.6), Inches(0.16), "是 ✓", size=7, bold=True, color=col_ok)

    # 否 → 向左到 Lane 2（不能处理分支）
    no_x  = d1x
    no_cy = L1_cy
    _conn(slide, no_x, no_cy, no_x - int(Inches(0.8)), no_cy, col_warn, width_pt=1.5)
    _label(slide, no_x - int(Inches(0.82)), no_cy - int(Inches(0.18)),
           Inches(0.7), Inches(0.16), "否 ✗", size=7, bold=True, color=col_warn)

    # ── Lane 2: 处理 & 存储 ────────────────────────────────────
    L2_top = ST + lane_h * 2
    L2_cy  = L2_top + lane_h // 2

    # ── 右侧分支：本地直接处理 ──
    rw_x = yes_x + int(Inches(0.8))
    rw_y = L2_cy - NH // 2
    r_ok = _rect(slide, rw_x, rw_y, NW, NH, col_ok, col_ok)
    add_text(r_ok, "本地处理\n完成", size=9, bold=True, color=col_white, align=PP_ALIGN.CENTER)

    # 连线 Lane1 → Lane2（是）
    _conn(slide, yes_x + int(Inches(0.8)), yes_cy, rw_x + NW // 2 - int(Inches(0.4)), L2_top, col_ok)
    _conn(slide, rw_x + NW // 2 - int(Inches(0.4)), L2_top,
          rw_x + NW // 2 - int(Inches(0.4)), rw_y, col_ok)

    # ── 左侧分支：不能处理 → 判断素材合格性 ──
    lw_x0 = no_x - int(Inches(0.8)) - NW  # 素材合格性判断菱形起点x
    lw_dx = lw_x0 - DW // 2 + NW // 2
    # 先画判断素材合格性节点（在 Lane 1 左侧区域）
    qa_x = no_x - int(Inches(0.9)) - DW
    qa_y = L1_cy - DH // 2
    qa = _diamond(slide, qa_x, qa_y, DW, DH, RGBColor(0xFF, 0xF5, 0xE6), col_warn)
    _label(slide, qa_x + int(Inches(0.1)), qa_y + int(Inches(0.14)),
           DW - int(Inches(0.2)), int(Inches(0.3)),
           "素材\n合格？", size=8, bold=True, color=col_warn)

    _conn(slide, no_x - int(Inches(0.8)), no_cy, qa_x + DW, L1_cy, col_warn)

    # 合格 → 存本地（等网络恢复）
    store_x = qa_x - int(Inches(0.14)) - NW
    store_y = L2_cy - NH // 2

    # 连线：合格 ↓
    qa_bottom_x = qa_x + DW // 2
    _conn(slide, qa_bottom_x, qa_y + DH, qa_bottom_x, L2_top, col_store)
    _conn(slide, qa_bottom_x, L2_top, qa_bottom_x, store_y, col_store)
    _label(slide, qa_x + int(Inches(0.1)), qa_y + DH + int(Inches(0.02)),
           DW - int(Inches(0.2)), Inches(0.14), "合格 ✓", size=7, bold=True, color=col_store)

    s_store = _rect(slide, store_x, store_y, NW, NH, col_store, col_store)
    add_text(s_store, "本地存储\n（等网络恢复）", size=8, bold=True, color=col_white, align=PP_ALIGN.CENTER)
    # 对齐 store X
    _conn(slide, qa_bottom_x, store_y + NH // 2,
          store_x + NW, store_y + NH // 2, col_store)

    # 不合格 → 丢弃/标记
    discard_x = qa_x + DW // 2 - NW // 2
    discard_y = L2_cy + int(Inches(0.12))
    discard_top_cx = qa_x + DW // 2
    # 向下箭头 - 不合格
    qa_top_x = qa_x + DW // 2
    # Actually draw the "不合格" branch going left-down
    d_box_x = qa_x - int(Inches(0.1)) - NW
    d_box_y = store_y + NH + int(Inches(0.12))
    if d_box_y < SB - int(Inches(0.05)):
        d_box = _rect(slide, d_box_x, d_box_y, NW, NH, RGBColor(0xC0, 0x39, 0x2B),
                      RGBColor(0xC0, 0x39, 0x2B))
        add_text(d_box, "标记不合格\n等待人工处理", size=8, bold=True, color=col_white, align=PP_ALIGN.CENTER)
        qa_left_x = qa_x
        _conn(slide, qa_left_x, L1_cy, d_box_x + NW, d_box_y + NH // 2, RGBColor(0xC0, 0x39, 0x2B))
        _label(slide, qa_left_x - int(Inches(0.85)), L1_cy + int(Inches(0.04)),
               Inches(0.75), Inches(0.16), "不合格 ✗", size=7, bold=True,
               color=RGBColor(0xC0, 0x39, 0x2B))

    # ── 结束事件：合并 ──────────────────────────────────────────
    end_x  = int(SL + SW - Inches(0.5)) - EV
    end_y  = L2_cy - EV // 2
    _oval(slide, end_x, end_y, EV, col_dark, col_dark)
    _conn(slide, r_ok_cx := rw_x + NW, L2_cy, end_x, L2_cy + EV // 2, col_dark)

    # ── 图例 ──────────────────────────────────────────────────
    leg_x = int(SL + SW - Inches(2.1))
    leg_y = ST + int(Inches(0.06))
    legends = [
        (col_primary, "任务流"),
        (col_warn,    "断网/异常"),
        (col_ok,      "本地处理"),
        (col_store,   "本地存储"),
    ]
    for i, (lc_item, lt) in enumerate(legends):
        lx = leg_x + i * int(Inches(0.52))
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, _i(lx), _i(leg_y + int(Inches(0.02))),
                                     _i(int(Inches(0.1))), _i(int(Inches(0.1))))
        dot.fill.solid()
        dot.fill.fore_color.rgb = lc_item
        dot.line.fill.background()
        _label(slide, lx + int(Inches(0.12)), leg_y, Inches(0.38), Inches(0.14),
               lt, size=6, color=col_dark, align=PP_ALIGN.LEFT)

    return slide


def main():
    def slides(ctx):
        build(ctx)

    prs = build_pptx(str(TEMPLATE), str(OUTPUT), slides)
    print(f"Output: {OUTPUT}")

    # QA
    from pptx import Presentation
    p = Presentation(str(OUTPUT))
    print(f"Slides: {len(p.slides)}")
    overflow = sum(
        1 for sl in p.slides
        for sh in sl.shapes
        if sh.left + sh.width > int(Inches(10.1))
    )
    print(f"Right overflow shapes: {overflow}")


if __name__ == "__main__":
    main()
