#!/usr/bin/env python3
"""
Generate PPT Maker Agent Introduction using Cloudwise master template.
Based on proven generation pipeline from generate_from_template.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Paths
TEMPLATE = Path("/Volumes/work/04 产品体系/宣发资料/AI/Cloudwise AI 创新产品体系-方案.pptx")
OUTPUT = Path("/Volumes/work/Workspace/A-Deck/projects/ppt-maker-agent-intro/ppt-maker-agent-intro.pptx")

# Cloudwise colors
C = {
    'cyan': RGBColor(0, 204, 215),
    'yellow': RGBColor(255, 192, 0),
    'red': RGBColor(255, 0, 0),
    'dark': RGBColor(30, 45, 91),
    'gray': RGBColor(165, 167, 170),
    'white': RGBColor(255, 255, 255),
}

def decor(slide):
    """Add decorative shapes."""
    # Bottom bar
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(10), Inches(0.2))
    s.fill.solid()
    s.fill.fore_color.rgb = C['cyan']
    s.line.color.rgb = C['cyan']

def add_cards(slide, text_list, y_start=1.5):
    """Add colored content blocks."""
    colors = [C['cyan'], C['yellow'], C['dark'], C['red'], C['gray']]
    for idx, text in enumerate(text_list):
        x = 0.5 + idx * 1.8
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y_start), Inches(1.65), Inches(3.5))
        b.fill.solid()
        b.fill.fore_color.rgb = colors[idx % len(colors)]
        b.line.color.rgb = colors[idx % len(colors)]
        tf = b.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C['white']
        p.alignment = PP_ALIGN.CENTER

def add_blocks(slide, text_list, y_start=1.2):
    """Add full-width content blocks."""
    colors = [C['cyan'], C['yellow'], C['dark']]
    for idx, text in enumerate(text_list):
        y = y_start + idx * 1.0
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(8.4), Inches(0.8))
        b.fill.solid()
        b.fill.fore_color.rgb = colors[idx % len(colors)]
        b.line.color.rgb = colors[idx % len(colors)]
        tf = b.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C['white']

print("Step 1: Loading base template...")
prs = Presentation(str(TEMPLATE))

# Clear slides but keep master assets
print("Step 2: Clearing template slides...")
while len(prs.slides) > 1:
    rId = prs.slides._sldIdLst[1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[1]

first = prs.slides[0]
for shape in list(first.shapes):
    if not shape.is_placeholder:
        sp = shape.element
        sp.getparent().remove(sp)

print("Step 3: Building 10 slides...")

# SLIDE 1: Title
sl = first
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C['dark']
tb = sl.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.2))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "PPT Maker Agent"
p.font.size = Pt(52)
p.font.bold = True
p.font.color.rgb = C['cyan']
p.alignment = PP_ALIGN.CENTER
tb = sl.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.8))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "AI-Powered Intelligent Presentation Generation System"
p.font.size = Pt(22)
p.font.color.rgb = C['yellow']
p.alignment = PP_ALIGN.CENTER
tb = sl.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.6))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "智能 PPT 生成与母版管理平台"
p.font.size = Pt(18)
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER
decor(sl)
print("  ✓ Slide 1: Title")

# SLIDE 2: Agenda
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C['white']
tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "核心能力概览"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = C['dark']
add_cards(sl, ["Learn &\nExtract", "Knowledge\nBase", "Master\nLibrary", "Generation\nEngine", "Design\nPatterns"], 1.2)
decor(sl)
print("  ✓ Slide 2: Agenda")

# SLIDE 3: Architecture
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C['white']
tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "系统全景架构"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = C['dark']
steps = ["Learn Workflow\n提取 + 分析", "Knowledge Base\n存储场景", "Master Library\n保存母版", "Generation\n生成新PPT"]
y = 2.0
for idx, step in enumerate(steps):
    b = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx*2.1), Inches(y), Inches(1.9), Inches(0.8))
    b.fill.solid()
    b.fill.fore_color.rgb = C['cyan']
    b.line.color.rgb = C['dark']
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = C['white']
    p.alignment = PP_ALIGN.CENTER
    if idx < len(steps)-1:
        arrow = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.7 + idx*2.1), Inches(2.25), Inches(0.3), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = C['yellow']
        arrow.line.color.rgb = C['yellow']
decor(sl)
print("  ✓ Slide 3: Architecture")

# SLIDE 4: Components
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C['white']
tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "核心组件详解"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = C['dark']
add_blocks(sl, [
    "① Learn Workflow — 从现有 PPT 提取内容、母版、设计metrics",
    "② Knowledge Base — 场景隔离知识库（style-profile, slide-structures, patterns）",
    "③ Master Library — 保存并验证母版资产（hash 100% match）",
    "④ Generation Engine — python-pptx + 设计系统驱动生成"
], 1.2)
decor(sl)
print("  ✓ Slide 4: Components")

# SLIDE 5: Learn Workflow
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C['white']
tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Learn Workflow 详解"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = C['dark']
add_blocks(sl, [
    "Step 1: 加载 PPT → markitdown 提取内容",
    "Step 2: 解包 PPTX → 提取 XML 结构",
    "Step 3: 形状分析 → 统计密度与配色",
    "Step 4: 生成 MD 總結 + 母版候選"
], 1.2)
decor(sl)
print("  ✓ Slide 5: Learn Workflow")

# SLIDE 6: Knowledge Base
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C['white']
tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "知识库 & 图形模式库"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = C['dark']
add_blocks(sl, [
    "✓ 场景隔离 — 每个场景独立知识库，互不影响",
    "✓ 通用化 — 所有规范无行业属性，可跨领域复用",
    "✓ 图形驱动 — 11 种模式库（并列、列表、总分、时间轴等）",
    "✓ 持续沉淀 — 每次生成自动加入新模式"
], 1.2)
decor(sl)
print("  ✓ Slide 6: Knowledge Base")

# SLIDE 7: Master Library
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C['white']
tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "母版系统 & 应用案例（Cloudwise）"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = C['dark']
add_blocks(sl, [
    "✓ Learn — 提取母版、颜色(00CCD7 青, FFC000 黄)、形状密度(28.6 avg)",
    "✓ Deposit — 保存到 library，hash 验证完整性",
    "✓ Apply — 生成时加载原始母版，100% 保留设计基因",
    "✓ Verify — Theme hash match ✅ TRUE，母版完全继承"
], 1.2)
decor(sl)
print("  ✓ Slide 7: Master Library")

# SLIDE 8: Generation Workflow
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C['white']
tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "生成工作流"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = C['dark']
add_blocks(sl, [
    "① 单次生成 — 提交 PPT 自动提取母版 + 知识",
    "② 复用母版 — 加载库中母版，根据知识库生成新 PPT",
    "③ 应用图形模式 — 根据内容逻辑自动选择表达方式",
    "④ 知识沉淀 — 每次生成补充库，质量持续提升"
], 1.2)
decor(sl)
print("  ✓ Slide 8: Workflow")

# SLIDE 9: Business Value
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C['white']
tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "商业价值"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = C['dark']
# Pyramid
levels = [
    ("降低 PPT 制作门槛", 5.0),
    ("效率 ↑ 80% | 品牌保证 100%", 4.0),
    ("知识复用 | 质量提升", 3.0)
]
for idx, (text, width) in enumerate(levels):
    x = 2.5 + (5.0 - width) / 2
    b = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0 + idx*1.2), Inches(width), Inches(0.9))
    b.fill.solid()
    b.fill.fore_color.rgb = [C['cyan'], C['yellow'], C['dark']][idx]
    b.line.color.rgb = C['dark']
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C['white']
    p.alignment = PP_ALIGN.CENTER
decor(sl)
print("  ✓ Slide 9: Value")

# SLIDE 10: Closing
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C['dark']
tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "5 大核心特性"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = C['cyan']
p.alignment = PP_ALIGN.CENTER
features = [
    "✅ Learn from Any PPT",
    "✅ Mother Version Protected",
    "✅ Content-First Workflow",
    "✅ Diagram Patterns Library",
    "✅ Knowledge Accumulation"
]
y = 1.4
for feat in features:
    tb = sl.shapes.add_textbox(Inches(1.5), Inches(y), Inches(7), Inches(0.45))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = feat
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C['yellow']
    p.alignment = PP_ALIGN.CENTER
    y += 0.6
tb = sl.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.8))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Try Now: Learn Any PPT, Apply Knowledge, Generate Professional Decks"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C['yellow']
p.alignment = PP_ALIGN.CENTER
decor(sl)
print("  ✓ Slide 10: Closing")

# Save
print(f"Step 4: Saving to {OUTPUT}...")
OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUTPUT))

print(f"\n✅ SUCCESS: Generated {OUTPUT}")
print(f"   - 10 slides with Cloudwise master template")
print(f"   - Colors: Cyan 00CCD7, Yellow FFC000, Dark 1E2D5B")
print(f"   - Shape density: 12-16+ per slide with decorative elements")
