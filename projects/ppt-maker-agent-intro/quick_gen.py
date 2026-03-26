#!/usr/bin/env python3
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

TEMPLATE = "/Volumes/work/04 产品体系/宣发资料/AI/Cloudwise AI 创新产品体系-方案.pptx"
OUTPUT = "/Volumes/work/Workspace/A-Deck/projects/ppt-maker-agent-intro/ppt-maker-agent-intro-v2.pptx"

C = {
    'cyan': RGBColor(0, 204, 215),
    'yellow': RGBColor(255, 192, 0),
    'dark': RGBColor(30, 45, 91),
    'white': RGBColor(255, 255, 255),
}

print("Loading...")
prs = Presentation(TEMPLATE)

# Clear
while len(prs.slides) > 1:
    rId = prs.slides._sldIdLst[1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[1]

first = prs.slides[0]
for shape in list(first.shapes):
    if not shape.is_placeholder:
        sp = shape.element
        sp.getparent().remove(sp)

# Slide 1 - Title (improved layout)
sl = first
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = C['dark']

title_box = sl.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "PPT Maker Agent"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = C['cyan']
p.alignment = PP_ALIGN.CENTER

subtitle_box = sl.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(0.6))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "AI-Powered Intelligent Presentation Generation System"
p.font.size = Pt(20)
p.font.color.rgb = C['yellow']
p.alignment = PP_ALIGN.CENTER

bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.9), Inches(10), Inches(0.15))
bar.fill.solid()
bar.fill.fore_color.rgb = C['cyan']
bar.line.color.rgb = C['cyan']

print("✓ Slide 1")

# Slide 2 - Agenda
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

title = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "核心能力概览"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = C['cyan']

items = ["Learn", "Knowledge", "Master", "Engine", "Patterns"]
colors = [C['cyan'], C['yellow'], C['dark'], RGBColor(255, 0, 0), RGBColor(165, 167, 170)]

for idx, item in enumerate(items):
    x = 0.5 + idx * 1.8
    card = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(1.65), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = colors[idx]
    card.line.width = Pt(0)
    
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C['white']
    p.alignment = PP_ALIGN.CENTER

bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.9), Inches(10), Inches(0.15))
bar.fill.solid()
bar.fill.fore_color.rgb = C['cyan']
print("✓ Slide 2")

# Slide 3-10: Content slides
slide_data = [
    ("系统全景架构", ["Learn Workflow", "Knowledge Base", "Master Library", "Generation"]),
    ("核心组件详解", ["① Learn", "② Knowledge", "③ Master", "④ Engine"]),
    ("Learn Workflow", ["Step1: 加载PPT", "Step2: 解包PPTX", "Step3: 形状分析", "Step4: 生成MD"]),
    ("知识库 & 图形模式", ["✓ 场景隔离", "✓ 通用化", "✓ 图形驱动", "✓ 持续沉淀"]),
    ("母版系统 & Cloudwise", ["✓ Learn", "✓ Deposit", "✓ Apply", "✓ Verify"]),
    ("生成工作流", ["单次生成", "复用母版", "应用图形模式", "知识沉淀"]),
    ("商业价值", ["降低门槛", "提高效率", "保证品牌", "提升质量"]),
    ("核心特性", ["✅ Learn from Any PPT", "✅ Mother Version", "✅ Content-First", "✅ Patterns & Knowledge"]),
]

for slide_idx, (title_text, items) in enumerate(slide_data, 3):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = RGBColor(255, 255, 255) if slide_idx < 10 else C['dark']
    
    # Title
    title_box = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C['cyan'] if slide_idx < 10 else C['cyan']
    
    # Content blocks
    colors_cycle = [C['cyan'], C['yellow'], C['dark'], RGBColor(255, 0, 0)]
    
    for item_idx, item in enumerate(items):
        y = 1.3 + item_idx * 1.0
        
        if slide_idx < 10:
            box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(8.4), Inches(0.85))
            box.fill.solid()
            box.fill.fore_color.rgb = colors_cycle[item_idx % len(colors_cycle)]
            box.line.width = Pt(0)
            
            tf = box.text_frame
            tf.margin_top = Inches(0.1)
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = C['white']
            p.alignment = PP_ALIGN.LEFT
        else:
            # Last slide - centered
            text_box = sl.shapes.add_textbox(Inches(1.0), Inches(y), Inches(8), Inches(0.5))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = C['yellow']
            p.alignment = PP_ALIGN.CENTER
    
    if slide_idx < 10:
        bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.9), Inches(10), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C['cyan']
    
    print(f"✓ Slide {slide_idx}")

Path(OUTPUT).parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT)

print(f"\n✅ Generated: {OUTPUT}")
