#!/usr/bin/env python3
"""Generate PPT Maker Agent self-introduction deck using Cloudwise master.

Design intent:
- Cloudwise master/theme preserved from template
- Product-brochure knowledge patterns: hierarchy, parallels, flow, list
- CMDB-aligned light style: white base, cyan emphasis, wide content coverage
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent / "skills" / "pptx" / "scripts"))
from pptx_lib import *  # noqa: E402,F403

from pptx.util import Inches
from pptx.enum.text import PP_ALIGN


PROJECT_DIR = Path(__file__).resolve().parent
TEMPLATE_PATH = PROJECT_DIR / "ppt-maker-agent-self-intro-cloudwise.pptx"
OUTPUT_PATH = PROJECT_DIR / "ppt-maker-agent-self-intro-cloudwise.pptx"


def _cover(slide):
    """Cover slide text blocks."""
    textbox(slide, Inches(0.72), Inches(1.26), Inches(12.11), Inches(1.00),
            "PPT Maker Agent", size=30, bold=True, color=C["cyan"], align=PP_ALIGN.CENTER)
    textbox(slide, Inches(0.80), Inches(2.42), Inches(12.03), Inches(0.70),
            "基于知识库与母版库的企业级 PPT 自动生成助手",
            size=22, bold=True, color=C["dark"], align=PP_ALIGN.CENTER)
    textbox(slide, Inches(0.96), Inches(3.30), Inches(11.76), Inches(0.80),
            "让 AI 按品牌规范、结构逻辑与业务语境自动生成可交付演示文稿",
            size=16, color=C["gray"], align=PP_ALIGN.CENTER)
    textbox(slide, Inches(1.09), Inches(4.80), Inches(11.49), Inches(0.45),
            "Cloudwise 母版继承  |  product-brochure 图形模式  |  方案型演示自动生成",
            size=14, color=C["dark"], align=PP_ALIGN.CENTER)
    textbox(slide, Inches(0.99), Inches(6.22), Inches(11.79), Inches(0.26),
            "Knowledge-driven  ·  Master-based  ·  Outline-first workflow",
            size=11, color=C["mid_gray"], align=PP_ALIGN.CENTER)


def _slide_1(prs, layout):
    slide = add_slide(prs, layout)
    _cover(slide)


def _slide_2(prs, layout):
    slide = add_slide(prs, layout, keep_title=True)
    header(slide, "PPT Maker Agent 是什么", "What it is", 2)
    render_relationship(slide, "我不是单纯的\"文字转 PPT\"工具", [
        {"title": "产品定位", "items": [
            "企业级 PPT 自动生成助手",
            "服务售前、产品、架构与汇报场景",
            "目标是可交付，而不是仅可预览",
        ]},
        {"title": "核心组成", "items": [
            "知识库",
            "母版库",
            "生成引擎",
        ]},
        {"title": "价值导向", "items": [
            "品牌一致性",
            "结构表达能力",
            "业务语境理解",
        ]},
    ])


def _slide_3(prs, layout):
    slide = add_slide(prs, layout, keep_title=True)
    header(slide, "当前能力概览", "Learning / Generation / QA", 3)
    render_three_cols(slide, [
        {"title": "学习能力", "items": [
            "从现有 PPT 提取配色、字体、布局规律",
            "识别品牌母版、主题与版式习惯",
            "将经验沉淀到 knowledge 持续复用",
            "支持参考材料风格学习与校准",
        ]},
        {"title": "生成能力", "items": [
            "基于 Cloudwise 母版输出新 PPT",
            "按内容逻辑匹配图形表达模式",
            "支持产品介绍、方案型、架构型页面",
            "支持提纲确认后的一次性生成",
        ]},
        {"title": "质检能力", "items": [
            "校验主题与母版继承是否正确",
            "抽取文本检查内容完整性",
            "避免空页、偏版与样式明显漂移",
            "支持生成后 QA 回读验证",
        ]},
    ])


def _slide_4(prs, layout):
    slide = add_slide(prs, layout, keep_title=True)
    header(slide, "PPT Maker Agent 工作原理", "Input → Knowledge → Master → Engine → Output", 4)
    render_flow(slide, [
        {"title": "输入需求", "desc": "主题 / 受众 / 页数 / 风格"},
        {"title": "Knowledge Base", "desc": "结构与图形表达规则"},
        {"title": "Master Library", "desc": "品牌主题、版式与样式资产"},
        {"title": "Generation Engine", "desc": "组装页面与编排内容"},
        {"title": "输出 PPT", "desc": "可交付的品牌化演示文稿"},
    ], y=FLOW_Y.inches)
    render_wide_list(slide, [
        "用户给出主题后，系统先确定内容结构，而不是直接画页面。",
        "Knowledge Base 决定“讲什么”和“怎么讲”；Master Library 决定“长什么样”。",
        "Generation Engine 负责把结构、样式与内容真正组装成一份 PPT。",
    ], [C["soft_blue"], C["light"], C["soft_yellow"]], y_start=LIST_Y.inches, h=0.64, gap=0.12, size=14)
    footer(slide)


def _slide_5(prs, layout):
    slide = add_slide(prs, layout, keep_title=True)
    header(slide, "我如何学习一个品牌风格", "Learn Workflow", 5)
    render_flow(slide, [
        {"title": "输入参考 PPT", "desc": "品牌资料 / 售前方案 / 产品稿"},
        {"title": "解包解析", "desc": "提取 XML / 版式 / theme"},
        {"title": "风格提炼", "desc": "色彩 / 字体 / 布局 / 图形规律"},
        {"title": "写入 Knowledge", "desc": "沉淀到 style / structures / patterns"},
        {"title": "后续复用", "desc": "在新生成任务中复用"},
    ], y=FLOW_Y.inches)
    render_wide_list(slide, [
        "可从 Cloudwise AI 宣发材料中学习通用结构。",
        "可从 CMDB 售前版中学习更贴近售前方案的样式。",
        "学习的目标不是复制某一页，而是复制“组织一页”的方法。",
    ], [C["light"], C["soft_blue"], C["soft_yellow"]], y_start=LIST_Y.inches, h=0.68, gap=0.12, size=14)


def _slide_6(prs, layout):
    slide = add_slide(prs, layout, keep_title=True)
    header(slide, "我如何生成一套 PPT", "Generation Workflow", 6)
    render_flow(slide, [
        {"title": "确认提纲", "desc": "先确认逐页内容"},
        {"title": "选择母版", "desc": "确定品牌模板与风格基线"},
        {"title": "匹配图形模式", "desc": "为每页选关系 / 流程 / 并列等"},
        {"title": "生成与填充", "desc": "落版式、写内容、套样式"},
        {"title": "QA 输出", "desc": "校验后交付 PPT"},
    ], y=FLOW_Y.inches)
    render_wide_list(slide, [
        "“先确认提纲，再生成”是固定流程，不再跳过。",
        "先把结构做对，再把样式做准，最后才是导出文件。",
        "这使生成结果更可控，也更适合真实业务交付。",
    ], [C["soft_yellow"], C["soft_blue"], C["light"]], y_start=LIST_Y.inches, h=0.68, gap=0.12, size=14)


def _slide_7(prs, layout):
    slide = add_slide(prs, layout, keep_title=True)
    header(slide, "我的知识库里有什么", "Knowledge Base Assets", 7)
    render_wide_list(slide, [
        "style-profile：色彩、字体、版式密度、风格标签与页面占位规则。",
        "diagram-patterns：并列、列表、总分、流程、关系等图形表达模式。",
        "slide-structures：产品介绍、售前方案、收尾页等结构模板。",
        "session-log：学习来源、提取记录与经验沉淀。",
        "UPDATES：风格校准与生成流程变更记录。",
    ], [C["soft_blue"], C["light"], C["soft_yellow"]], y_start=1.35, h=0.86, gap=0.14, size=15)


def _slide_8(prs, layout):
    slide = add_slide(prs, layout, keep_title=True)
    header(slide, "我的母版能力", "Master Library Capabilities", 8)
    render_three_cols(slide, [
        {"title": "品牌继承", "items": [
            "沿用 Cloudwise 的主题、颜色和版心逻辑",
            "确保输出在视觉上属于同一品牌体系",
            "减少风格漂移和局部自定义失控",
        ]},
        {"title": "版式复用", "items": [
            "复用 slide master / layout / theme",
            "让不同类型页面共享统一结构基础",
            "支持后续多套母版并行扩展",
        ]},
        {"title": "输出一致性", "items": [
            "验证主题 hash 与母版继承",
            "保证字体、色板、布局连续一致",
            "让演示稿具备正式交付质量",
        ]},
    ])


def _slide_9(prs, layout):
    slide = add_slide(prs, layout, keep_title=True)
    header(slide, "我的业务价值", "Faster / Safer / Reusable", 9)
    render_three_cols(slide, [
        {"title": "更快", "items": [
            "从需求到提纲再到初稿，缩短准备周期",
            "减少手工排版与反复调样式成本",
            "支持快速迭代和多版本输出",
        ]},
        {"title": "更稳", "items": [
            "减少品牌风格漂移与版式偏差",
            "避免内容堆砌与页面稀疏失衡",
            "输出结果更适合真实售前和汇报场景",
        ]},
        {"title": "更可复制", "items": [
            "同一套知识可服务多个产品线",
            "母版与知识沉淀可以持续复用",
            "适合方案、产品、架构和汇报等多场景",
        ]},
    ])
    footer(slide)


def _slide_10(prs, layout):
    slide = add_slide(prs, layout)
    rrect(slide, SAFE_LEFT, Inches(3.84), SAFE_WIDTH, Inches(1.24), C["soft_blue"], line=C["line"])
    textbox(slide, SAFE_LEFT + Inches(0.2), Inches(1.45), SAFE_WIDTH - Inches(0.24), Inches(0.9),
             "PPT 制作，让 AI 更像一名售前顾问",
             size=28, bold=True, color=C["dark"], align=PP_ALIGN.CENTER)
    textbox(slide, SAFE_LEFT + Inches(0.42), Inches(2.55), SAFE_WIDTH - Inches(0.72), Inches(0.9),
             "我不只是生成页面，我在学习品牌、理解结构、组织表达、输出结果。",
             size=16, color=C["gray"], align=PP_ALIGN.CENTER)
    textbox(slide, SAFE_LEFT + Inches(0.36), Inches(4.22), SAFE_WIDTH - Inches(0.62), Inches(0.5),
             "给我一个主题、一套母版和一个目标受众，我会先给你提纲确认，再交付 PPT 初稿。",
             size=15, bold=True, color=C["dark"], align=PP_ALIGN.CENTER)
    stats = [
        ("Cloudwise Master", "品牌母版继承"),
        ("Knowledge Driven", "知识库驱动生成"),
        ("Outline First", "先确认提纲再出稿"),
    ]
    left = SAFE_LEFT
    gap = Inches(0.18)
    width = SAFE_WIDTH
    box_w = (width - gap * 2) / 3
    for idx, (title, desc) in enumerate(stats):
        x = left + idx * (box_w + gap)
        rrect(slide, x, Inches(5.45), box_w, Inches(0.95), C["light"], line=C["line"])
        textbox(slide, x + Inches(0.04), Inches(5.62), box_w - Inches(0.08), Inches(0.25),
                 title, size=12, bold=True, color=C["cyan"], align=PP_ALIGN.CENTER)
        textbox(slide, x + Inches(0.04), Inches(5.93), box_w - Inches(0.08), Inches(0.2),
                 desc, size=10, color=C["gray"], align=PP_ALIGN.CENTER)
    footer(slide)


def build() -> Path:
    prs = Presentation(str(TEMPLATE_PATH))
    scale_x = prs.slide_width / Inches(BASE_WIDTH_IN)
    clear_slides(prs)
    cover_layout = layout_by_names(prs, ["封面"], 0)
    thanks_layout = layout_by_names(prs, ["感谢页"], 7)
    content_layout = layout_by_names(prs, ["內容页", "内容页"], 1)

    _slide_1(prs, cover_layout)
    fit_slide_to_canvas(prs.slides[-1], scale_x)
    _slide_2(prs, content_layout)
    fit_slide_to_canvas(prs.slides[-1], scale_x)
    _slide_3(prs, content_layout)
    fit_slide_to_canvas(prs.slides[-1], scale_x)
    _slide_4(prs, content_layout)
    fit_slide_to_canvas(prs.slides[-1], scale_x)
    _slide_5(prs, content_layout)
    fit_slide_to_canvas(prs.slides[-1], scale_x)
    _slide_6(prs, content_layout)
    fit_slide_to_canvas(prs.slides[-1], scale_x)
    _slide_7(prs, content_layout)
    fit_slide_to_canvas(prs.slides[-1], scale_x)
    _slide_8(prs, content_layout)
    fit_slide_to_canvas(prs.slides[-1], scale_x)
    _slide_9(prs, content_layout)
    fit_slide_to_canvas(prs.slides[-1], scale_x)
    _slide_10(prs, thanks_layout)
    fit_slide_to_canvas(prs.slides[-1], scale_x)

    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


if __name__ == "__main__":
    output = build()
    print(f"Generated: {output}")
