#!/usr/bin/env python3
"""Generate A-Deck self-intro PPT (v2) — 12 slides, three masters.

Pages
  1  产品定位全景  三列 (定位 / 痛点 Before-After / 范围)
  2  系统总览 L0  四层 (输入 / 双库 / 引擎5子模块 / 输出)
  3  业务流程 L1  五阶段横向流
  4  L2 泳道图 A+B  Layout Mapping + Data Binding
  5  L2 泳道图 C+D  Theme Resolution + Geometry Scaling
  6  L2 泳道图 E+F  QA Gates + Regression & Release
  7  能力地图  L1→L2→产出物 三列格子
  8  机制一  Master + Spec 协作
  9  机制二  Skill 协作
 10  机制三  质量门禁与证据链
 11  多母版回归实证  矩阵
 12  落地路径 + 收尾
"""

from pathlib import Path
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent.parent
SCRIPTS_DIR = ROOT / "skills" / "pptx" / "scripts"
SHARED_DIR  = ROOT / "skills" / "pptx" / "shared"
TOGAF_LOADERS_DIR = ROOT / "skills" / "pptx" / "togaf-architecture" / "loaders"
XSEAN_LOADERS_DIR = ROOT / "skills" / "pptx" / "XSean" / "loaders"
sys.path.insert(0, str(SCRIPTS_DIR))
sys.path.insert(0, str(SHARED_DIR))
sys.path.insert(0, str(TOGAF_LOADERS_DIR))   # real TOGAF loaders
sys.path.append(str(XSEAN_LOADERS_DIR))

from pptx_lib import build_pptx, header, verify_pptx  # noqa: E402
from renderer_utils import configure_theme, shape_rect, textbox  # noqa: E402
# P7 / P9 delegate directly to the real TOGAF render functions
from common import (  # noqa: E402
    render_capability_grid,
    render_flow_stages,
    render_journey_stages,
    render_matrix_table,
    render_swimlane,
)
from x01_function_architecture import load_slide as load_x01_function_architecture  # noqa: E402

PROJECT_DIR = Path(__file__).resolve().parent

MASTERS = [
    ("dark-cloudwise-green", "self-intro-dark-cloudwise-green.pptx"),
    ("light-cloudwise-cyan", "self-intro-light-cloudwise-cyan.pptx"),
    ("light-cloudwise-purple", "self-intro-light-cloudwise-purple.pptx"),
]

# ─── Low-level helpers ────────────────────────────────────────────────────────

def set_subtitle(ctx, slide, text):
    ph = next(
        (s for s in slide.shapes
         if getattr(s, "is_placeholder", False) and s.placeholder_format.idx == 1),
        None,
    )
    if ph and ph.has_text_frame:
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(ctx.typography["subtitle"])
        p.font.color.rgb = RGBColor(*ctx.colors["gray"])


def _lighter(c, a=110):
    return tuple(min(255, v + a) for v in c)


def _arrow_h(slide, x, mid_y, col):
    """Horizontal right-pointing arrow: line + solid arrowhead block."""
    aw = int(Inches(0.14))
    shape_rect(slide, x, mid_y - int(Inches(0.018)),
               aw - int(Inches(0.04)), int(Inches(0.036)), fill_color=col)
    shape_rect(slide, x + aw - int(Inches(0.05)), mid_y - int(Inches(0.06)),
               int(Inches(0.05)), int(Inches(0.12)), fill_color=col)


def _lane_palette(C):
    return [
        C["primary"],
        C["secondary"],
        C.get("accent5", C["primary"]),
        C.get("accent6", C["secondary"]),
    ]


def render_loader_content_slide(ctx, title, subtitle, render_fn, data, *, top_pad=0.22, bottom_pad=0.1):
    """Render a loader into the current master's safe content region."""
    slide = ctx.add_content_slide()
    header(slide, title)
    set_subtitle(ctx, slide, subtitle)

    x = int(ctx.safe_left)
    y = int(ctx.content_top + Inches(top_pad))
    w = int(ctx.safe_width)
    h = int(ctx.safe_bottom - y - Inches(bottom_pad))
    render_fn(ctx, data, "A-Deck", slide=slide, region=(x, y, w, h))
    return slide


# ─── Page 1: 产品定位全景（三列） ────────────────────────────────────────────

def add_positioning_overview(ctx):
    slide = ctx.add_content_slide()
    C = ctx.colors
    header(slide, "A-Deck · 产品定位全景")
    set_subtitle(ctx, slide, "定位 / 痛点 / 范围  三要素一览")

    x = int(ctx.safe_left)
    y = int(ctx.content_top + Inches(0.2))
    w = int(ctx.safe_width)
    h = int(ctx.safe_bottom - y - Inches(0.12))
    gap = int(Inches(0.16))
    cw = int((w - gap * 2) / 3)

    # ── 左列：产品定位 ──
    lx = x
    shape_rect(slide, lx, y, cw, h, fill_color=C["dark"], line_color=C["primary"])
    textbox(slide, lx + int(Inches(0.1)), y + int(Inches(0.12)),
            cw - int(Inches(0.2)), int(Inches(0.26)),
            "A-Deck 是什么",
            size="label", bold=True, color=C["white"])
    textbox(slide, lx + int(Inches(0.1)), y + int(Inches(0.46)),
            cw - int(Inches(0.2)), int(Inches(0.55)),
            "企业级可交付 PPT\n生成与编辑智能体",
            size="caption", bold=True, color=C["white"])
    for i, tag in enumerate(["Knowledge-driven", "Master-based", "QA-gated"]):
        ty = y + int(Inches(1.16 + i * 0.4))
        shape_rect(slide, lx + int(Inches(0.1)), ty,
                   cw - int(Inches(0.2)), int(Inches(0.3)),
                   fill_color=C["primary"], line_color=C["line"])
        textbox(slide, lx + int(Inches(0.12)), ty + int(Inches(0.07)),
                cw - int(Inches(0.24)), int(Inches(0.18)),
                tag, size="micro", bold=True, color=C["white"])
    textbox(slide, lx + int(Inches(0.1)), y + int(Inches(2.5)),
            cw - int(Inches(0.2)), int(Inches(0.22)),
            "应用场景：售前  ·  产品  ·  架构  ·  运营",
            size="micro", color=C["white"])

    # ── 中列：痛点 Before / After ──
    mx = x + cw + gap
    shape_rect(slide, mx, y, cw, h, fill_color=C["light"], line_color=C["line"])
    textbox(slide, mx + int(Inches(0.1)), y + int(Inches(0.1)),
            cw - int(Inches(0.2)), int(Inches(0.24)),
            "手工制作  →  智能产线",
            size="label", bold=True, color=C["text"])
    asis_h = int(h * 0.44)
    ay = y + int(Inches(0.42))
    shape_rect(slide, mx + int(Inches(0.08)), ay,
               cw - int(Inches(0.16)), asis_h,
               fill_color=C.get("accent1", (200, 210, 220)),
               line_color=C["line"])
    textbox(slide, mx + int(Inches(0.12)), ay + int(Inches(0.06)),
            cw - int(Inches(0.24)), int(Inches(0.2)),
            "As-Is", size="micro", bold=True, color=C["text"])
    for i, t in enumerate(["手工排版，反复改版", "多人风格漂移不一", "质量依赖个人经验"]):
        textbox(slide, mx + int(Inches(0.14)), ay + int(Inches(0.3 + i * 0.28)),
                cw - int(Inches(0.28)), int(Inches(0.22)),
                f"• {t}", size="micro", color=C["text"])
    tobe_y = ay + asis_h + int(Inches(0.1))
    tobe_h = h - (tobe_y - y) - int(Inches(0.1))
    shape_rect(slide, mx + int(Inches(0.08)), tobe_y,
               cw - int(Inches(0.16)), tobe_h,
               fill_color=C.get("accent4", (210, 240, 220)), line_color=C["primary"])
    textbox(slide, mx + int(Inches(0.12)), tobe_y + int(Inches(0.06)),
            cw - int(Inches(0.24)), int(Inches(0.2)),
            "To-Be  A-Deck", size="micro", bold=True, color=C["dark"])
    for i, t in enumerate(["提纲先行，先确认后生成", "母版继承，品牌一致", "自动质检，可交付输出"]):
        textbox(slide, mx + int(Inches(0.14)), tobe_y + int(Inches(0.32 + i * 0.28)),
                cw - int(Inches(0.28)), int(Inches(0.22)),
                f"• {t}", size="micro", color=C["dark"])
    textbox(slide, mx + int(Inches(0.1)), tobe_y + int(Inches(1.26)),
            cw - int(Inches(0.2)), int(Inches(0.2)),
            "→ 交付效率与一致性显著提升",
            size="micro", bold=True, color=C["dark"])

    # ── 右列：范围边界 ──
    rx = mx + cw + gap
    shape_rect(slide, rx, y, cw, h, fill_color=C["light"], line_color=C["line"])
    textbox(slide, rx + int(Inches(0.1)), y + int(Inches(0.1)),
            cw - int(Inches(0.2)), int(Inches(0.24)),
            "范围与目标",
            size="label", bold=True, color=C["text"])
    sections = [
        ("Scope In",  ["W1-W6 全流程能力", "路线图/架构/产品PPT", "多母版跨主题生成"],  C["primary"]),
        ("Scope Out", ["非PPT格式文档",     "离线演示草稿",        "非结构化自由创作"], C.get("accent1", (200, 210, 220))),
        ("质量目标",  ["可交付优先（非草稿）", "QA证据链齐全",      "回归测试可复现"],  C["secondary"]),
    ]
    each_h = int((h - int(Inches(0.42))) / 3)
    sy = y + int(Inches(0.42))
    for label, items, col in sections:
        shape_rect(slide, rx + int(Inches(0.08)), sy,
                   cw - int(Inches(0.16)), int(Inches(0.26)), fill_color=col)
        textbox(slide, rx + int(Inches(0.1)), sy + int(Inches(0.05)),
                cw - int(Inches(0.2)), int(Inches(0.18)),
                label, size="micro", bold=True, color=C["white"])
        for i, t in enumerate(items):
            textbox(slide, rx + int(Inches(0.12)), sy + int(Inches(0.32 + i * 0.26)),
                    cw - int(Inches(0.24)), int(Inches(0.22)),
                    f"• {t}", size="micro", color=C["text"])
        sy += each_h + int(Inches(0.06))


# ─── Page 2: 系统总览 L0（四层，引擎细化 5 子模块） ──────────────────────────

def add_l0_overview(ctx):
    slide = ctx.add_content_slide()
    C = ctx.colors
    header(slide, "系统总览（L0）")
    set_subtitle(ctx, slide, "输入  →  双库  →  生成引擎（5子系统）  →  可交付输出")

    x = int(ctx.safe_left)
    w = int(ctx.safe_width)
    cy = int(ctx.content_top + Inches(0.1))
    L_GAP = int(Inches(0.2))

    # ─ Layer 1: 输入 ─
    LH = int(Inches(0.58))
    shape_rect(slide, x, cy, w, LH, fill_color=C["light"], line_color=C["line"])
    textbox(slide, x + int(Inches(0.12)), cy + int(Inches(0.17)),
            int(Inches(0.9)), int(Inches(0.22)),
            "输入", size="label", bold=True, color=C["text"])
    inputs = ["主题 / 受众", "时长 / 深度", "风格约束", "场景标签"]
    iw = int((w - int(Inches(1.1)) - int(Inches(0.12)) * 3) / 4)
    ix0 = x + int(Inches(1.1))
    for i, lb in enumerate(inputs):
        bx = ix0 + i * (iw + int(Inches(0.12)))
        shape_rect(slide, bx, cy + int(Inches(0.1)), iw, int(Inches(0.38)),
                   fill_color=C["primary"], line_color=C["line"])
        textbox(slide, bx, cy + int(Inches(0.18)), iw, int(Inches(0.22)),
                lb, size="micro", bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    cy += LH + L_GAP

    # ─ Layer 2: 双库 ─
    LH = int(Inches(1.12))
    shape_rect(slide, x, cy, w, LH,
               fill_color=C.get("accent4", (210, 240, 220)), line_color=C["line"])
    half_w = int((w - int(Inches(0.16))) / 2)
    kb_x = x + int(Inches(0.08))
    shape_rect(slide, kb_x, cy + int(Inches(0.08)),
               half_w - int(Inches(0.08)), LH - int(Inches(0.16)),
               fill_color=C["secondary"], line_color=C["line"])
    textbox(slide, kb_x + int(Inches(0.1)), cy + int(Inches(0.13)),
            half_w - int(Inches(0.28)), int(Inches(0.22)),
            "Knowledge Base", size="label", bold=True, color=C["white"])
    for i, lbl in enumerate(["产品 Spec 文档库", "图形 Skill 库", "Loader 装载器库"]):
        textbox(slide, kb_x + int(Inches(0.12)), cy + int(Inches(0.43 + i * 0.22)),
                half_w - int(Inches(0.32)), int(Inches(0.2)),
                f"• {lbl}", size="micro", color=C["white"])
    ml_x = x + int(Inches(0.08)) + half_w + int(Inches(0.08))
    shape_rect(slide, ml_x, cy + int(Inches(0.08)),
               half_w - int(Inches(0.08)), LH - int(Inches(0.16)),
               fill_color=C["dark"], line_color=C["primary"])
    textbox(slide, ml_x + int(Inches(0.1)), cy + int(Inches(0.13)),
            half_w - int(Inches(0.28)), int(Inches(0.22)),
            "Master Library", size="label", bold=True, color=C["white"])
    for i, lbl in enumerate(["dark / cyan / purple 三套母版",
                               "master.pptx + cloudwise-spec.yaml",
                               "theme XML → palette → ctx 注入"]):
        textbox(slide, ml_x + int(Inches(0.12)), cy + int(Inches(0.43 + i * 0.22)),
                half_w - int(Inches(0.32)), int(Inches(0.2)),
                f"• {lbl}", size="micro", color=C["white"])
    cy += LH + L_GAP

    # ─ Layer 3: 生成引擎（5 子模块） ─
    LH = int(Inches(1.0))
    shape_rect(slide, x, cy, w, LH, fill_color=C["primary"], line_color=C["line"])
    textbox(slide, x + int(Inches(0.12)), cy + int(Inches(0.12)),
            int(Inches(1.3)), int(Inches(0.22)),
            "生成引擎", size="label", bold=True, color=C["white"])
    textbox(slide, x + int(Inches(0.12)), cy + int(Inches(0.38)),
            int(Inches(1.3)), int(Inches(0.45)),
            "Orchestrator\n编排与分发\n↕ 5 子模块", size="micro", color=C["white"])
    engines = [
        ("Layout Mapping",   "版式精确匹配\n+ fallback 兜底"),
        ("Loader Registry",  "按slide_id调度\nloader(ctx,data)"),
        ("Theme Resolution", "theme XML解析\n→ ctx.colors注入"),
        ("Scaling Adapter",  "EMU整数化\n防双重缩放"),
    ]
    avail_w = w - int(Inches(1.5)) - int(Inches(0.1)) * 3
    ew = int(avail_w / 4)
    ex0 = x + int(Inches(1.5))
    for i, (etitle, esub) in enumerate(engines):
        bx = ex0 + i * (ew + int(Inches(0.1)))
        col = C["dark"] if i % 2 == 0 else C["secondary"]
        shape_rect(slide, bx, cy + int(Inches(0.1)), ew, LH - int(Inches(0.2)),
                   fill_color=col, line_color=C["line"])
        textbox(slide, bx + int(Inches(0.03)), cy + int(Inches(0.14)),
                ew - int(Inches(0.06)), int(Inches(0.22)),
                etitle, size="micro", bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        textbox(slide, bx + int(Inches(0.03)), cy + int(Inches(0.4)),
                ew - int(Inches(0.06)), int(Inches(0.42)),
                esub, size="micro", color=C["white"], align=PP_ALIGN.CENTER)
    cy += LH + L_GAP

    # ─ Layer 4: 输出 ─
    LH = int(Inches(0.58))
    shape_rect(slide, x, cy, w, LH, fill_color=C["light"], line_color=C["line"])
    textbox(slide, x + int(Inches(0.12)), cy + int(Inches(0.17)),
            int(Inches(0.9)), int(Inches(0.22)),
            "输出", size="label", bold=True, color=C["text"])
    outputs = ["PPTX 交付文件", "theme match 报告", "overflow 检测", "slides 计数验证"]
    ow = int((w - int(Inches(1.1)) - int(Inches(0.12)) * 3) / 4)
    ox0 = x + int(Inches(1.1))
    for i, lb in enumerate(outputs):
        bx = ox0 + i * (ow + int(Inches(0.12)))
        shape_rect(slide, bx, cy + int(Inches(0.1)), ow, int(Inches(0.38)),
                   fill_color=C["secondary"], line_color=C["line"])
        textbox(slide, bx, cy + int(Inches(0.17)), ow, int(Inches(0.22)),
                lb, size="micro", bold=True, color=C["white"], align=PP_ALIGN.CENTER)


# ─── Page 3: 业务流程 L1（5 阶段） ───────────────────────────────────────────

def add_l1_flow(ctx):
    data = {
        "title": "业务流程（L1）",
        "subtitle": "从需求到交付的五阶段端到端流程",
        "content": {
            "stages": [
                {"title": "需求澄清", "items": ["确认目标受众", "时长与交付形式"]},
                {"title": "提纲确认", "items": ["先确认提纲", "再进入生成"]},
                {"title": "母版+映射", "items": ["选择母版", "内容与版式匹配"]},
                {"title": "生成编排", "items": ["Loader 调度", "PPTX 渲染输出"]},
                {"title": "QA+交付", "items": ["三层门禁验证", "输出 QA 证据"]},
            ],
            "note": "BA-02 Stage Flow | A-Deck 端到端流程 | render_flow_stages",
        },
    }
    return render_loader_content_slide(
        ctx,
        data["title"],
        data["subtitle"],
        render_flow_stages,
        data,
        top_pad=0.28,
        bottom_pad=0.14,
    )


# ─── Pages 4-6: L2 泳道图（TOGAF swimlane 样式） ─────────────────────────────

def draw_swimlane(ctx, slide, x, y, w, h, title, lanes, C):
    """Compatibility wrapper: keep local half-slide title bar, delegate body to loader."""
    title_h = int(Inches(0.28))
    shape_rect(slide, x, y, w, title_h, fill_color=C["dark"], line_color=C["line"])
    textbox(slide, x + int(Inches(0.06)), y + int(Inches(0.05)),
            w - int(Inches(0.12)), int(Inches(0.2)),
            title, size="caption", bold=True, color=C["white"], align=PP_ALIGN.CENTER)

    region_top = y + title_h + int(Inches(0.04))
    region_height = h - title_h - int(Inches(0.04))
    data = {
        "content": {"lanes": lanes},
    }
    return render_swimlane(
        ctx,
        data,
        "A-Deck",
        slide=slide,
        region=(x, region_top, w, region_height),
    )


def add_l2_swimlane_slide(ctx, page_no, subtitle, left_title, left_lanes,
                           right_title, right_lanes):
    slide = ctx.add_content_slide()
    C = ctx.colors
    header(slide, f"L2 执行流程（{page_no}）")
    set_subtitle(ctx, slide, subtitle)

    x = int(ctx.safe_left)
    y = int(ctx.content_top + Inches(0.22))
    w = int(ctx.safe_width)
    h = int(ctx.safe_bottom - y - Inches(0.1))
    gap = int(Inches(0.18))
    lw = int((w - gap) / 2)

    draw_swimlane(ctx, slide, x, y, lw, h, left_title, left_lanes, C)
    draw_swimlane(ctx, slide, x + lw + gap, y, lw, h, right_title, right_lanes, C)


# ─── Page 7: X01 功能架构图 ───────────────────────────────────────────────────

def add_capability_map(ctx):
    """Render P7 as X01 function architecture diagram."""
    data = {
        "title": "A-Deck 功能架构总览",
        "subtitle": "X01 · 分层 > 分组 > 模块 的汇报型架构表达",
        "content": {
            "side_labels": {
                "left": "生成链路",
                "right": "交付约束",
            },
            "layers": [
                {
                    "id": "context",
                    "name": "目标 / 上下文层",
                    "groups": [
                        {
                            "name": "输入场景",
                            "modules": [
                                {"id": "need", "name": "需求输入"},
                                {"id": "outline", "name": "提纲确认"},
                            ],
                        },
                        {
                            "name": "交付目标",
                            "modules": [
                                {"id": "brand", "name": "品牌一致"},
                                {"id": "deliver", "name": "可交付PPT"},
                            ],
                        },
                    ],
                },
                {
                    "id": "service",
                    "name": "能力 / 服务层",
                    "groups": [
                        {
                            "name": "Learn",
                            "modules": [
                                {"id": "learn_style", "name": "样式提取"},
                                {"id": "learn_struct", "name": "结构抽取"},
                                {"id": "learn_rule", "name": "规则沉淀"},
                            ],
                        },
                        {
                            "name": "Generate",
                            "modules": [
                                {"id": "gen_content", "name": "内容编排"},
                                {"id": "gen_layout", "name": "版式映射"},
                                {"id": "gen_graphic", "name": "图形装配"},
                            ],
                        },
                        {
                            "name": "QA / Edit",
                            "modules": [
                                {"id": "qa_gate", "name": "质量门禁"},
                                {"id": "xml_fix", "name": "XML精修"},
                                {"id": "partial_regen", "name": "局部重生成"},
                            ],
                        },
                    ],
                },
                {
                    "id": "core",
                    "name": "核心对象层",
                    "groups": [
                        {
                            "name": "知识与母版",
                            "modules": [
                                {"id": "kb", "name": "Knowledge Base", "type": "storage"},
                                {"id": "skill", "name": "Skill Library", "type": "storage"},
                                {"id": "master", "name": "Master Library", "type": "storage"},
                            ],
                        },
                        {
                            "name": "运行时对象",
                            "modules": [
                                {"id": "registry", "name": "Loader Registry"},
                                {"id": "ctx_build", "name": "BuildContext"},
                                {"id": "evidence", "name": "QA Evidence", "type": "storage"},
                            ],
                        },
                    ],
                },
                {
                    "id": "input",
                    "name": "输入 / 来源层",
                    "groups": [
                        {
                            "name": "结构化输入",
                            "modules": [
                                {"id": "yaml", "name": "Config YAML"},
                                {"id": "master_pptx", "name": "Master PPTX"},
                                {"id": "spec_doc", "name": "Spec / Loader Spec"},
                            ],
                        },
                    ],
                },
                {
                    "id": "output",
                    "name": "输出 / 承载层",
                    "groups": [
                        {
                            "name": "交付产物",
                            "modules": [
                                {"id": "pptx_out", "name": "PPTX产物"},
                                {"id": "report_out", "name": "验证报告"},
                                {"id": "matrix_out", "name": "回归矩阵"},
                            ],
                        },
                    ],
                },
            ],
            "external_zones": [
                {
                    "title": "外部对象 / 对接区",
                    "style": "dashed",
                    "items": ["用户需求", "企业素材", "品牌规范", "交付场景"],
                },
            ],
            "links": [
                {"from": "input", "to": "service", "label": "驱动"},
                {"from": "core", "to": "service", "label": "支撑"},
                {"from": "service", "to": "output", "label": "输出"},
                {"from": "service", "to": "context", "label": "达成"},
            ],
        },
    }
    return load_x01_function_architecture(ctx, data)


# ─── Page 8: 机制一  Master + Spec 协作 ──────────────────────────────────────

def add_master_spec_mechanism(ctx):
    slide = ctx.add_content_slide()
    C = ctx.colors
    header(slide, "机制一：Master + Spec 协作")
    set_subtitle(ctx, slide, "theme XML 是主题真源，spec 是约束与兼容层")

    x = int(ctx.safe_left)
    y = int(ctx.content_top + Inches(0.28))
    w = int(ctx.safe_width)
    gap = int(Inches(0.2))
    bw = int((w - gap * 3) / 4)
    bh = int(Inches(1.5))

    comps = [
        ("Master PPTX",         ["布局占位符定义", "母版视觉样式", "theme1.xml（真源）"], False),
        ("cloudwise-spec.yaml", ["palette 别名映射", "typography 配置", "safe_zone 约束"], True),
        ("BuildContext",        ["ctx.colors 注入", "ctx.palette 注入", "ctx.typography 下发"], False),
        ("Loader",              ["apply fill/line/text", "禁止品牌色硬编码", "只读ctx，不读配置"], True),
    ]
    for i, (title, items, dark) in enumerate(comps):
        bx = x + i * (bw + gap)
        bg = C["dark"] if dark else C["light"]
        tc = C["white"] if dark else C["text"]
        shape_rect(slide, bx, y, bw, bh, fill_color=bg, line_color=C["line"])
        textbox(slide, bx + int(Inches(0.08)), y + int(Inches(0.1)),
                bw - int(Inches(0.16)), int(Inches(0.28)),
                title, size="label", bold=True, color=tc, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            textbox(slide, bx + int(Inches(0.1)), y + int(Inches(0.5 + j * 0.3)),
                    bw - int(Inches(0.2)), int(Inches(0.24)),
                    f"• {item}", size="micro", color=tc)
        if i < 3:
            _arrow_h(slide, bx + bw + int(Inches(0.04)), y + bh // 2, C["primary"])

    cy = y + bh + int(Inches(0.28))
    shape_rect(slide, x, cy, w, int(Inches(0.38)), fill_color=C["primary"])
    textbox(slide, x + int(Inches(0.2)), cy + int(Inches(0.09)),
            w - int(Inches(0.4)), int(Inches(0.22)),
            "结论：theme XML → spec 合并 → ctx 下发 → loader 渲染，全链路无品牌色硬编码",
            size="caption", bold=True, color=C["white"], align=PP_ALIGN.CENTER)


# ─── Page 9: PPT生成场景旅程 — 直接调用 render_journey_stages（BA-08 真实 loader）──

def add_scenario_journey(ctx):
    """Delegate to the real TOGAF BA-08 render_journey_stages loader."""
    data = {
        "title": "PPT 生成旅程图",
        "subtitle": "BA-08 · 从需求输入到 QA 交付的五阶段场景旅程",
        "content": {
            "phases": [
                {"title": "需求输入",
                 "touchpoints": ["Config YAML 规范化输入", "约束校验"],
                 "emotion": "~",
                 "pain": "无结构化输入，遗漏关键约束"},
                {"title": "提纲确认",
                 "touchpoints": ["Outline Gate 先行", "待确认后生成"],
                 "emotion": "+",
                 "pain": "直接生成，事后反复大改"},
                {"title": "图形编排",
                 "touchpoints": ["Layout Mapping", "Loader Registry 调度"],
                 "emotion": "+",
                 "pain": "图形风格不一致"},
                {"title": "渲染生成",
                 "touchpoints": ["Theme Resolution", "Scaling Adapter"],
                 "emotion": "+",
                 "pain": "颜色硬编码，换母版全改"},
                {"title": "QA 交付",
                 "touchpoints": ["三层门禁验证", "QA 证据输出"],
                 "emotion": "+",
                 "pain": ""},
            ],
            "note": "BA-08 Scenario Journey | A-Deck PPT生成场景 | render_journey_stages",
        },
    }
    return render_loader_content_slide(
        ctx,
        data["title"],
        data["subtitle"],
        render_journey_stages,
        data,
        top_pad=0.22,
        bottom_pad=0.1,
    )


# ─── Page 10: 机制三  质量门禁与证据链 ───────────────────────────────────────

def add_qa_gates(ctx):
    slide = ctx.add_content_slide()
    C = ctx.colors
    header(slide, "机制三：质量门禁与证据链")
    set_subtitle(ctx, slide, "三层 QA · 不通过不交付 · 全程可追溯")

    x = int(ctx.safe_left)
    y = int(ctx.content_top + Inches(0.28))
    w = int(ctx.safe_width)
    rh = int(Inches(0.72))
    rgap = int(Inches(0.14))

    qa_layers = [
        ("内容 QA",  ["文本完整性", "必填字段覆盖", "关键词一致性"], 92, C["primary"]),
        ("视觉 QA",  ["视觉一致性", "theme色对齐",  "overflow检测"], 88, C["secondary"]),
        ("结构 QA",  ["theme hash匹配", "页数符合预期", "占位符完整"], 96, C["dark"]),
    ]
    for i, (name, checks, score, col) in enumerate(qa_layers):
        ry = y + i * (rh + rgap)
        shape_rect(slide, x, ry, w, rh, fill_color=C["light"], line_color=C["line"])
        shape_rect(slide, x, ry, int(Inches(0.22)), rh, fill_color=col)
        textbox(slide, x + int(Inches(0.3)), ry + int(Inches(0.18)),
                int(Inches(1.4)), int(Inches(0.24)),
                name, size="label", bold=True, color=C["text"])
        for j, ck in enumerate(checks):
            textbox(slide, x + int(Inches(1.9 + j * 1.9)), ry + int(Inches(0.2)),
                    int(Inches(1.7)), int(Inches(0.22)),
                    f"✓ {ck}", size="micro", color=C["text"])
        bar_x = x + int(Inches(7.6))
        bar_total = int(w - int(Inches(7.8)))
        bar_w = int(bar_total * score / 100)
        shape_rect(slide, bar_x, ry + int(Inches(0.2)),
                   bar_total, int(Inches(0.24)),
                   fill_color=C.get("accent6", (220, 220, 220)))
        shape_rect(slide, bar_x, ry + int(Inches(0.2)), bar_w, int(Inches(0.24)), fill_color=col)
        textbox(slide, x + w - int(Inches(0.7)), ry + int(Inches(0.2)),
                int(Inches(0.6)), int(Inches(0.24)),
                f"{score}%", size="caption", bold=True, color=C["text"], align=PP_ALIGN.RIGHT)

    ey = y + len(qa_layers) * (rh + rgap) + int(Inches(0.12))
    shape_rect(slide, x, ey, w, int(Inches(0.5)), fill_color=C["dark"], line_color=C["primary"])
    textbox(slide, x + int(Inches(0.15)), ey + int(Inches(0.13)),
            w - int(Inches(0.3)), int(Inches(0.24)),
            "证据链：QA 摘要报告  ·  回归矩阵（slides / theme match / overflow）  ·  可交付版本标记",
            size="caption", bold=True, color=C["white"], align=PP_ALIGN.CENTER)


# ─── Page 11: 多母版回归实证 ──────────────────────────────────────────────────

def add_regression_matrix(ctx):
    data = {
        "title": "多母版回归实证",
        "subtitle": "dark / cyan / purple × GTM / Roadmap / TOGAF 全矩阵验证通过",
        "content": {
            "row_headers": [
                "Theme Match",
                "Overflow",
                "GTM Slides",
                "TOGAF Slides",
                "Roadmap Slides",
            ],
            "col_headers": ["Dark Green", "Cyan", "Purple"],
            "cells": [
                ["✓", "✓", "✓"],
                ["0", "0", "0"],
                ["35", "35", "35"],
                ["40", "40", "40"],
                ["12", "12", "12"],
            ],
            "note": "BA-06 Matrix Table | 多母版回归矩阵 | render_matrix_table",
        },
    }
    return render_loader_content_slide(
        ctx,
        data["title"],
        data["subtitle"],
        render_matrix_table,
        data,
        top_pad=0.22,
        bottom_pad=0.1,
    )


# ─── Page 12: 落地路径 + 收尾 ─────────────────────────────────────────────────

def add_implementation_path(ctx):
    slide = ctx.add_content_slide()
    C = ctx.colors
    header(slide, "落地路径与下一步合作")
    set_subtitle(ctx, slide, "你给场景与受众 → 我交提纲确认 → 生成初稿 → QA 证据")

    x = int(ctx.safe_left)
    y = int(ctx.content_top + Inches(0.22))
    w = int(ctx.safe_width)
    gap = int(Inches(0.2))
    bw = int((w - gap * 2) / 3)
    bh = int(Inches(2.5))

    phases = [
        ("1", "PoC",   "1–2 周", ["目标：验证单场景可行性",    "输入：1套母版 + 1个场景",   "输出：10页PPTX + QA报告",     "验收：theme match + 0 overflow"]),
        ("2", "Pilot", "2–4 周", ["目标：多场景试点运行",       "输入：2-3套母版 + N个场景", "输出：N个PPTX + 回归矩阵",    "验收：基线稳定 + 可复现"]),
        ("3", "Scale", "持续",   ["目标：资产化与常态运营",     "输入：Skill库 + 规则库沉淀","输出：自动化回归 + 存量维护", "验收：扩展新需求零改引擎"]),
    ]
    for i, (num, phase, duration, items) in enumerate(phases):
        bx = x + i * (bw + gap)
        shape_rect(slide, bx, y, bw, bh, fill_color=C["light"], line_color=C["line"])
        shape_rect(slide, bx, y, bw, int(Inches(0.44)), fill_color=C["primary"])
        textbox(slide, bx + int(Inches(0.08)), y + int(Inches(0.1)),
                int(Inches(0.26)), int(Inches(0.26)),
                num, size="label", bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        textbox(slide, bx + int(Inches(0.42)), y + int(Inches(0.1)),
                bw - int(Inches(0.52)), int(Inches(0.26)),
                f"{phase}  ·  {duration}", size="label", bold=True, color=C["white"])
        for j, item in enumerate(items):
            textbox(slide, bx + int(Inches(0.14)), y + int(Inches(0.62 + j * 0.42)),
                    bw - int(Inches(0.28)), int(Inches(0.34)),
                    f"• {item}", size="caption", color=C["text"])
        if i < 2:
            _arrow_h(slide, bx + bw + int(Inches(0.04)), y + bh // 2, C["primary"])

    cta_y = y + bh + int(Inches(0.24))
    shape_rect(slide, x, cta_y, w, int(Inches(0.46)), fill_color=C["dark"], line_color=C["primary"])
    textbox(slide, x + int(Inches(0.2)), cta_y + int(Inches(0.12)),
            w - int(Inches(0.4)), int(Inches(0.24)),
            "下一步：选 1 个场景进入 PoC  →  你给主题 + 受众  →  我交 QA 可交付版 PPTX",
            size="caption", bold=True, color=C["white"], align=PP_ALIGN.CENTER)


# ─── Orchestrator ─────────────────────────────────────────────────────────────

def my_slides(ctx):
    configure_theme(ctx.colors, ctx.template_spec)

    add_positioning_overview(ctx)   # P1 三列全景

    add_l0_overview(ctx)            # P2 L0 四层（引擎5子模块）

    add_l1_flow(ctx)                # P3 L1 五阶段

    # P4 L2-A Layout Mapping + L2-B Data Binding (detailed nodes+links)
    slide = ctx.add_content_slide()
    C = ctx.colors
    header(slide, "L2 执行流程（1/3）")
    set_subtitle(ctx, slide, "Layout Mapping（版式映射）+ Data Binding（数据装配）")
    
    x = int(ctx.safe_left)
    y = int(ctx.content_top + Inches(0.22))
    w = int(ctx.safe_width)
    h = int(ctx.safe_bottom - y - Inches(0.1))
    gap = int(Inches(0.18))
    lw = int((w - gap) / 2)
    
    # L2-A Layout Mapping (left half)
    title_h = int(Inches(0.28))
    shape_rect(slide, x, y, lw, title_h, fill_color=C["dark"], line_color=C["line"])
    textbox(slide, x + int(Inches(0.06)), y + int(Inches(0.05)),
            lw - int(Inches(0.12)), int(Inches(0.2)),
            "L2-A  Layout Mapping", size="caption", bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    
    render_swimlane(
        ctx,
        {
            "content": {
                "lanes": ["Orchestrator", "PPTX Layout", "Fallback", "Placeholder"],
                "columns": 5,
                "nodes": [
                    {"id": "start", "type": "event", "lane": "Orchestrator", "x_in": 0.02, "start": True},
                    {"id": "n1", "type": "step", "lane": "Orchestrator", "col": 0, "name": "读取layout索引", "system": "Config"},
                    {"id": "n2", "type": "step", "lane": "PPTX Layout", "col": 1, "name": "名称精确匹配", "system": "Layout"},
                    {"id": "d1", "type": "decision", "lane": "PPTX Layout", "col": 2, "text": "匹配成功?", "w_in": 0.48, "h_in": 0.42},
                    {"id": "n3", "type": "step", "lane": "Placeholder", "col": 3, "name": "检查占位符", "system": "Master"},
                    {"id": "n4", "type": "step", "lane": "Fallback", "col": 2, "name": "回退兜底", "system": "Index", "manual": True},
                    {"id": "end", "type": "event", "lane": "Placeholder", "x_in": 5.08, "end": True},
                ],
                "links": [
                    {"from": "start", "to": "n1"},
                    {"from": "n1", "to": "n2"},
                    {"from": "n2", "to": "d1"},
                    {"from": "d1", "to": "n3", "label": "Yes"},
                    {"from": "d1", "to": "n4", "label": "No", "dashed": True, "start_side": "bottom", "end_side": "top"},
                    {"from": "n3", "to": "end"},
                    {"from": "n4", "to": "end", "dashed": True},
                ],
            }
        },
        "A-Deck",
        slide=slide,
        region=(x, y + title_h + int(Inches(0.04)), lw, h - title_h - int(Inches(0.04))),
    )
    
    # L2-B Data Binding (right half)
    shape_rect(slide, x + lw + gap, y, lw, title_h, fill_color=C["dark"], line_color=C["line"])
    textbox(slide, x + lw + gap + int(Inches(0.06)), y + int(Inches(0.05)),
            lw - int(Inches(0.12)), int(Inches(0.2)),
            "L2-B  Data Binding", size="caption", bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    
    render_swimlane(
        ctx,
        {
            "content": {
                "lanes": ["Config YAML", "Validator", "Loader Registry", "Renderer"],
                "columns": 5,
                "nodes": [
                    {"id": "start", "type": "event", "lane": "Config YAML", "x_in": 0.02, "start": True},
                    {"id": "n1", "type": "step", "lane": "Config YAML", "col": 0, "name": "读取配置", "system": "File"},
                    {"id": "n2", "type": "step", "lane": "Validator", "col": 1, "name": "校验必填字段", "system": "Schema"},
                    {"id": "d1", "type": "decision", "lane": "Validator", "col": 2, "text": "有效?", "w_in": 0.42, "h_in": 0.42},
                    {"id": "n3", "type": "step", "lane": "Loader Registry", "col": 3, "name": "分发Loader", "system": "Registry"},
                    {"id": "n4", "type": "step", "lane": "Renderer", "col": 4, "name": "渲染输出", "system": "PPTX", "critical": True},
                    {"id": "end", "type": "event", "lane": "Renderer", "x_in": 5.08, "end": True},
                ],
                "links": [
                    {"from": "start", "to": "n1"},
                    {"from": "n1", "to": "n2"},
                    {"from": "n2", "to": "d1"},
                    {"from": "d1", "to": "n3", "label": "Yes"},
                    {"from": "n3", "to": "n4"},
                    {"from": "n4", "to": "end"},
                ],
            }
        },
        "A-Deck",
        slide=slide,
        region=(x + lw + gap, y + title_h + int(Inches(0.04)), lw, h - title_h - int(Inches(0.04))),
    )

    # P5 L2-C Theme Resolution + L2-D Geometry Scaling (detailed nodes+links)
    slide = ctx.add_content_slide()
    C = ctx.colors
    header(slide, "L2 执行流程（2/3）")
    set_subtitle(ctx, slide, "Theme Resolution（主题色解析）+ Geometry Scaling（几何缩放）")
    
    x = int(ctx.safe_left)
    y = int(ctx.content_top + Inches(0.22))
    w = int(ctx.safe_width)
    h = int(ctx.safe_bottom - y - Inches(0.1))
    gap = int(Inches(0.18))
    lw = int((w - gap) / 2)
    
    # L2-C Theme Resolution (left half)
    title_h = int(Inches(0.28))
    shape_rect(slide, x, y, lw, title_h, fill_color=C["dark"], line_color=C["line"])
    textbox(slide, x + int(Inches(0.06)), y + int(Inches(0.05)),
            lw - int(Inches(0.12)), int(Inches(0.2)),
            "L2-C  Theme Resolution", size="caption", bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    
    render_swimlane(
        ctx,
        {
            "content": {
                "lanes": ["Theme XML", "Spec YAML", "BuildContext", "Loader"],
                "columns": 5,
                "nodes": [
                    {"id": "start", "type": "event", "lane": "Theme XML", "x_in": 0.02, "start": True},
                    {"id": "n1", "type": "step", "lane": "Theme XML", "col": 0, "name": "解析theme1.xml", "system": "Parser"},
                    {"id": "n2", "type": "step", "lane": "Spec YAML", "col": 1, "name": "读palette别名", "system": "Spec"},
                    {"id": "n3", "type": "step", "lane": "BuildContext", "col": 2, "name": "合并ctx.colors", "system": "Merger", "critical": True},
                    {"id": "n4", "type": "step", "lane": "Loader", "col": 3, "name": "apply fill/line", "system": "Render"},
                    {"id": "end", "type": "event", "lane": "Loader", "x_in": 5.08, "end": True},
                ],
                "links": [
                    {"from": "start", "to": "n1"},
                    {"from": "n1", "to": "n2"},
                    {"from": "n2", "to": "n3"},
                    {"from": "n3", "to": "n4"},
                    {"from": "n4", "to": "end"},
                ],
            }
        },
        "A-Deck",
        slide=slide,
        region=(x, y + title_h + int(Inches(0.04)), lw, h - title_h - int(Inches(0.04))),
    )
    
    # L2-D Geometry Scaling (right half)
    shape_rect(slide, x + lw + gap, y, lw, title_h, fill_color=C["dark"], line_color=C["line"])
    textbox(slide, x + lw + gap + int(Inches(0.06)), y + int(Inches(0.05)),
            lw - int(Inches(0.12)), int(Inches(0.2)),
            "L2-D  Geometry Scaling", size="caption", bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    
    render_swimlane(
        ctx,
        {
            "content": {
                "lanes": ["Canvas Detector", "Scale Calculator", "EMU Converter", "QA Checker"],
                "columns": 5,
                "nodes": [
                    {"id": "start", "type": "event", "lane": "Canvas Detector", "x_in": 0.02, "start": True},
                    {"id": "n1", "type": "step", "lane": "Canvas Detector", "col": 0, "name": "读画布EMU", "system": "Reader"},
                    {"id": "n2", "type": "step", "lane": "Scale Calculator", "col": 1, "name": "计算缩放", "system": "Calc"},
                    {"id": "n3", "type": "step", "lane": "EMU Converter", "col": 2, "name": "坐标整数化", "system": "Convert", "critical": True},
                    {"id": "n4", "type": "step", "lane": "QA Checker", "col": 3, "name": "边界检测", "system": "Check"},
                    {"id": "n5", "type": "decision", "lane": "QA Checker", "col": 4, "text": "安全?", "w_in": 0.42, "h_in": 0.42},
                    {"id": "end", "type": "event", "lane": "QA Checker", "x_in": 5.08, "end": True},
                ],
                "links": [
                    {"from": "start", "to": "n1"},
                    {"from": "n1", "to": "n2"},
                    {"from": "n2", "to": "n3"},
                    {"from": "n3", "to": "n4"},
                    {"from": "n4", "to": "n5"},
                    {"from": "n5", "to": "end"},
                ],
            }
        },
        "A-Deck",
        slide=slide,
        region=(x + lw + gap, y + title_h + int(Inches(0.04)), lw, h - title_h - int(Inches(0.04))),
    )

    # P6 L2-E QA Gates + L2-F Regression & Release (detailed nodes+links)
    slide = ctx.add_content_slide()
    C = ctx.colors
    header(slide, "L2 执行流程（3/3）")
    set_subtitle(ctx, slide, "QA Gates（质量门禁）+ Regression & Release（回归发布）")
    
    x = int(ctx.safe_left)
    y = int(ctx.content_top + Inches(0.22))
    w = int(ctx.safe_width)
    h = int(ctx.safe_bottom - y - Inches(0.1))
    gap = int(Inches(0.18))
    lw = int((w - gap) / 2)
    
    # L2-E QA Gates (left half)
    title_h = int(Inches(0.28))
    shape_rect(slide, x, y, lw, title_h, fill_color=C["dark"], line_color=C["line"])
    textbox(slide, x + int(Inches(0.06)), y + int(Inches(0.05)),
            lw - int(Inches(0.12)), int(Inches(0.2)),
            "L2-E  QA Gates", size="caption", bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    
    render_swimlane(
        ctx,
        {
            "content": {
                "lanes": ["Content QA", "Visual QA", "Structure QA", "Release Gate"],
                "columns": 5,
                "nodes": [
                    {"id": "start", "type": "event", "lane": "Content QA", "x_in": 0.02, "start": True},
                    {"id": "n1", "type": "step", "lane": "Content QA", "col": 0, "name": "文本完整性", "system": "Parser"},
                    {"id": "n2", "type": "step", "lane": "Visual QA", "col": 1, "name": "视觉一致性", "system": "Validator"},
                    {"id": "n3", "type": "step", "lane": "Structure QA", "col": 2, "name": "theme hash校验", "system": "Hash"},
                    {"id": "d1", "type": "decision", "lane": "Structure QA", "col": 3, "text": "全部通过?", "w_in": 0.52, "h_in": 0.42},
                    {"id": "n4", "type": "step", "lane": "Release Gate", "col": 4, "name": "输出证据链", "system": "Report", "critical": True},
                    {"id": "end", "type": "event", "lane": "Release Gate", "x_in": 5.08, "end": True},
                ],
                "links": [
                    {"from": "start", "to": "n1"},
                    {"from": "n1", "to": "n2"},
                    {"from": "n2", "to": "n3"},
                    {"from": "n3", "to": "d1"},
                    {"from": "d1", "to": "n4", "label": "Yes"},
                    {"from": "n4", "to": "end"},
                ],
            }
        },
        "A-Deck",
        slide=slide,
        region=(x, y + title_h + int(Inches(0.04)), lw, h - title_h - int(Inches(0.04))),
    )
    
    # L2-F Regression & Release (right half)
    shape_rect(slide, x + lw + gap, y, lw, title_h, fill_color=C["dark"], line_color=C["line"])
    textbox(slide, x + lw + gap + int(Inches(0.06)), y + int(Inches(0.05)),
            lw - int(Inches(0.12)), int(Inches(0.2)),
            "L2-F  Regression & Release", size="caption", bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    
    render_swimlane(
        ctx,
        {
            "content": {
                "lanes": ["Generator", "QA Runner", "Baseline Comparator", "Release Manager"],
                "columns": 5,
                "nodes": [
                    {"id": "start", "type": "event", "lane": "Generator", "x_in": 0.02, "start": True},
                    {"id": "n1", "type": "step", "lane": "Generator", "col": 0, "name": "三母版生成", "system": "Multi", "critical": True},
                    {"id": "n2", "type": "step", "lane": "QA Runner", "col": 1, "name": "汇总QA结果", "system": "Runner"},
                    {"id": "n3", "type": "step", "lane": "Baseline Comparator", "col": 2, "name": "基线比对", "system": "Compare"},
                    {"id": "d1", "type": "decision", "lane": "Baseline Comparator", "col": 3, "text": "异常?", "w_in": 0.48, "h_in": 0.42},
                    {"id": "n4", "type": "step", "lane": "Release Manager", "col": 4, "name": "标记可交付", "system": "Release", "critical": True},
                    {"id": "end", "type": "event", "lane": "Release Manager", "x_in": 5.08, "end": True},
                ],
                "links": [
                    {"from": "start", "to": "n1"},
                    {"from": "n1", "to": "n2"},
                    {"from": "n2", "to": "n3"},
                    {"from": "n3", "to": "d1"},
                    {"from": "d1", "to": "n4", "label": "No"},
                    {"from": "n4", "to": "end"},
                ],
            }
        },
        "A-Deck",
        slide=slide,
        region=(x + lw + gap, y + title_h + int(Inches(0.04)), lw, h - title_h - int(Inches(0.04))),
    )

    add_capability_map(ctx)         # P7
    add_master_spec_mechanism(ctx)  # P8
    add_scenario_journey(ctx)        # P9 BA-08 Scenario Journey Map
    add_qa_gates(ctx)               # P10
    add_regression_matrix(ctx)      # P11
    add_implementation_path(ctx)    # P12


# ─── QA + Main ────────────────────────────────────────────────────────────────

def check_overflow(pptx_path):
    prs = Presentation(str(pptx_path))
    sw = prs.slide_width
    overflow = [
        i for i, sl in enumerate(prs.slides, 1)
        if max((sh.left + sh.width for sh in sl.shapes), default=0) > sw
    ]
    return len(prs.slides), overflow


def main():
    print("=" * 64)
    for master_name, out_name in MASTERS:
        template = ROOT / "skills" / "pptx" / "master-library" / master_name / "cloudwise-master.pptx"
        spec = ROOT / "skills" / "pptx" / "master-library" / master_name / "cloudwise-spec.yaml"
        output = PROJECT_DIR / out_name

        print(f"\n▶ Generating: {master_name}")
        build_pptx(template, output, my_slides, spec)
        verify_pptx(output, template)
        slides, overflow = check_overflow(output)
        print(f"  Right overflow: {len(overflow)}/{slides}")
        if overflow:
            print(f"  Overflow slides: {overflow}")
        print(f"  Output: {output}")

    print("\n✓ Done")


if __name__ == "__main__":
    main()
