#!/usr/bin/env python3
"""Generate one-slide Service Interaction diagram for PPT Maker Agent.

Uses Cloudwise master template and content layout (内容1/內容页).
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent / "skills" / "pptx" / "scripts"))
from pptx_lib import *  # noqa: E402,F403

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TEMPLATE = Path("/Volumes/work/Workspace/A-Deck/projects/ppt-maker-intro-cloudwise-2025"
                "/ppt-maker-agent-self-intro-cloudwise.pptx")
OUTPUT = Path(__file__).resolve().parent / "ppt-maker-service-interaction.pptx"


def _i(v):
    """Convert any computed coordinate to integer EMU."""
    return int(v)


def _node(slide, x, y, w, h, title, role,
          fill, text_color, border=None, dashed=False, bold=False):
    """Draw a service node with title and role subtitle."""
    shape = rrect(slide, x, y, w, h, fill, line=border or fill, adj=3500)
    shape.line.width = Pt(2.0 if bold else 1.2)
    if dashed:
        set_dash(shape)

    textbox(slide, x + Inches(0.08), y + Inches(0.12), w - Inches(0.16), Inches(0.28),
            title, size=10, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    textbox(slide, x + Inches(0.08), y + Inches(0.40), w - Inches(0.16), Inches(0.20),
            role, size=7, color=text_color, align=PP_ALIGN.CENTER)

    return shape


def _seq_badge(slide, x, y, n, color):
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        _i(x - Inches(0.08)), _i(y - Inches(0.08)), _i(Inches(0.16)), _i(Inches(0.16)),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    textbox(slide, x - Inches(0.06), y - Inches(0.065), Inches(0.12), Inches(0.12),
            str(n), size=8, bold=True, color=C["white"], align=PP_ALIGN.CENTER)


def _line(slide, x1, y1, x2, y2, color, width=1.8, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _i(x1), _i(y1), _i(x2), _i(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if dashed:
        conn.line.dash_style = 4
    return conn


def _arrow_head(slide, x, y, color, direction="right"):
    if direction == "right":
        tri = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            _i(x - Inches(0.06)), _i(y - Inches(0.045)), _i(Inches(0.09)), _i(Inches(0.09)),
        )
        tri.rotation = 90
    elif direction == "left":
        tri = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            _i(x - Inches(0.03)), _i(y - Inches(0.045)), _i(Inches(0.09)), _i(Inches(0.09)),
        )
        tri.rotation = 270
    elif direction == "down":
        tri = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            _i(x - Inches(0.045)), _i(y - Inches(0.03)), _i(Inches(0.09)), _i(Inches(0.09)),
        )
        tri.rotation = 180
    else:
        tri = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            _i(x - Inches(0.045)), _i(y - Inches(0.06)), _i(Inches(0.09)), _i(Inches(0.09)),
        )
        tri.rotation = 0
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()


def _label(slide, x, y, text, size=8, color=None, italic=False):
    box = textbox(slide, x, y, Inches(1.6), Inches(0.2), text,
                  size=size, color=color or C["gray"], align=PP_ALIGN.CENTER)
    if italic:
        box.text_frame.paragraphs[0].font.italic = True


def my_slides(ctx):
    slide = ctx.add_content_slide()
    header(slide, "PPT Maker Agent — Service Interaction")
    textbox(slide, SAFE_LEFT + Inches(0.04), Inches(0.68), Inches(6.4), Inches(0.24),
            "Scenario: Generate enterprise PPT from outline (Flow + Event)",
            size=10, color=C["brand_gray"])

    node_w = Inches(1.55)
    node_h = Inches(0.78)
    y_top = Inches(1.55)

    x1 = SAFE_LEFT + Inches(0.05)
    x2 = x1 + Inches(1.77)
    x3 = x2 + Inches(1.77)
    x4 = x3 + Inches(1.77)
    x5 = x4 + Inches(1.77)

    _node(slide, x1, y_top, node_w, node_h, "API Gateway", "Initiating Service",
          fill=C["domain_bg"], text_color=C["white"], border=C["domain_bg"])
    _node(slide, x2, y_top, node_w, node_h, "Generation Orchestrator", "Core Domain Service",
          fill=C["cyan"], text_color=C["white"], border=C["cyan"], bold=True)
    _node(slide, x3, y_top, node_w, node_h, "Knowledge Service", "Core Domain Service",
          fill=C["cyan"], text_color=C["white"], border=C["cyan"])
    _node(slide, x4, y_top, node_w, node_h, "Template Service", "Supporting Service",
          fill=C["cyan_2"], text_color=C["dark"], border=C["cyan"])
    _node(slide, x5, y_top, node_w, node_h, "LLM Provider", "External Service",
          fill=C["brand_gray"], text_color=C["white"], border=C["brand_gray"], dashed=True)

    channel_x = SAFE_LEFT + Inches(4.15)
    channel_y = Inches(3.88)
    channel = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, channel_x, channel_y, Inches(1.2), Inches(0.56))
    channel.fill.solid()
    channel.fill.fore_color.rgb = C["domain_bg"]
    channel.line.color.rgb = C["domain_bg"]
    textbox(slide, channel_x + Inches(0.08), channel_y + Inches(0.08), Inches(1.04), Inches(0.4),
            "Event Bus\nppt.generated", size=8, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

    storage_x = SAFE_LEFT + Inches(6.55)
    storage_y = Inches(3.86)
    _node(slide, storage_x, storage_y, Inches(1.75), node_h, "Object Storage", "Supporting Service",
          fill=C["cyan_2"], text_color=C["dark"], border=C["cyan"])

    y_mid = y_top + Inches(0.39)

    # 1: gateway -> orchestrator
    _line(slide, x1 + node_w, y_mid, x2, y_mid, C["cyan"])
    _arrow_head(slide, x2, y_mid, C["cyan"], "right")
    _seq_badge(slide, (x1 + node_w + x2) / 2, y_mid - Inches(0.11), 1, C["cyan"])
    _label(slide, x1 + Inches(1.56), y_mid - Inches(0.29), "POST /generate", color=C["dark"])
    _label(slide, x1 + Inches(1.56), y_mid - Inches(0.13), "HTTP", size=7, color=C["brand_gray"], italic=True)

    # 2: orchestrator -> knowledge
    _line(slide, x2 + node_w, y_mid, x3, y_mid, C["cyan"])
    _arrow_head(slide, x3, y_mid, C["cyan"], "right")
    _seq_badge(slide, (x2 + node_w + x3) / 2, y_mid - Inches(0.11), 2, C["cyan"])
    _label(slide, x2 + Inches(1.56), y_mid - Inches(0.29), "Context Lookup (<=200ms)", color=C["dark"])
    _label(slide, x2 + Inches(1.56), y_mid - Inches(0.13), "gRPC", size=7, color=C["brand_gray"], italic=True)

    # 3: orchestrator -> template
    _line(slide, x2 + node_w, y_mid + Inches(0.24), x4, y_mid + Inches(0.24), C["cyan_2"])
    _arrow_head(slide, x4, y_mid + Inches(0.24), C["cyan_2"], "right")
    _seq_badge(slide, (x2 + node_w + x4) / 2, y_mid + Inches(0.13), 3, C["cyan_2"])
    _label(slide, x3 + Inches(0.42), y_mid + Inches(0.02), "Layout Resolve", color=C["dark"])

    # 4: orchestrator -> llm
    _line(slide, x2 + node_w, y_mid + Inches(0.48), x5, y_mid + Inches(0.48), C["cyan"])
    _arrow_head(slide, x5, y_mid + Inches(0.48), C["cyan"], "right")
    _seq_badge(slide, (x2 + node_w + x5) / 2, y_mid + Inches(0.37), 4, C["cyan"])
    _label(slide, x4 + Inches(0.2), y_mid + Inches(0.28), "Prompt Completion", color=C["dark"])
    _label(slide, x4 + Inches(0.2), y_mid + Inches(0.44), "HTTPS", size=7, color=C["brand_gray"], italic=True)

    # 5: response llm -> orchestrator
    _line(slide, x5, y_mid + Inches(0.66), x2 + node_w, y_mid + Inches(0.66), C["brand_gray"], width=1.1)
    _arrow_head(slide, x2 + node_w, y_mid + Inches(0.66), C["brand_gray"], "left")
    _seq_badge(slide, (x2 + node_w + x5) / 2, y_mid + Inches(0.77), 5, C["brand_gray"])
    _label(slide, x4 + Inches(0.12), y_mid + Inches(0.78), "Draft JSON Response", color=C["gray"])

    # 6: orchestrator -> event bus (async)
    o_x = x2 + node_w / 2
    o_y = y_top + node_h
    _line(slide, o_x, o_y, channel_x + Inches(0.55), channel_y, C["domain_bg"], dashed=True)
    _arrow_head(slide, channel_x + Inches(0.55), channel_y, C["domain_bg"], "down")
    _seq_badge(slide, (o_x + channel_x + Inches(0.55)) / 2, (o_y + channel_y) / 2 - Inches(0.08), 6, C["domain_bg"])
    _label(slide, SAFE_LEFT + Inches(2.7), Inches(3.1), "Publish deck.generated", color=C["dark"])
    _label(slide, SAFE_LEFT + Inches(2.7), Inches(3.25), "Kafka", size=7, color=C["brand_gray"], italic=True)

    # 7: event bus -> storage (async)
    _line(slide, channel_x + Inches(1.2), channel_y + Inches(0.28), storage_x, storage_y + Inches(0.39), C["domain_bg"], dashed=True)
    _arrow_head(slide, storage_x, storage_y + Inches(0.39), C["domain_bg"], "right")
    _seq_badge(slide, channel_x + Inches(1.64), channel_y + Inches(0.15), 7, C["domain_bg"])
    _label(slide, channel_x + Inches(1.15), channel_y + Inches(0.37), "Store deck asset", color=C["dark"])

    # Failure path: llm timeout -> orchestrator fallback
    _line(slide, x5 - Inches(0.06), y_top + Inches(0.08), x2 + node_w + Inches(0.1), y_top + Inches(0.08),
          RGBColor(0xC0, 0x39, 0x2B), width=1.1, dashed=True)
    _arrow_head(slide, x2 + node_w + Inches(0.1), y_top + Inches(0.08), RGBColor(0xC0, 0x39, 0x2B), "left")
    _label(slide, x4 - Inches(0.2), y_top - Inches(0.16), "Error Path: timeout -> fallback template", size=7,
           color=RGBColor(0xC0, 0x39, 0x2B))

    textbox(slide, SAFE_LEFT, Inches(7.12), SAFE_WIDTH, Inches(0.2),
            "Contract: API v2.1  |  Retry: 3 attempts with exponential backoff  |  Critical SLA: <=200ms",
            size=7, color=C["brand_gray"], align=PP_ALIGN.LEFT)


if __name__ == "__main__":
    build_pptx(TEMPLATE, OUTPUT, my_slides)
    verify_pptx(OUTPUT, TEMPLATE)
