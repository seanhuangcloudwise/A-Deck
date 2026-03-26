#!/usr/bin/env python3
"""Test TOGAF generation on all 3 masters, each to a separate output file."""
import re, sys
from pathlib import Path

import yaml
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parent.parent.parent
SCRIPTS_DIR = ROOT / "skills" / "pptx" / "scripts"
LOADERS_DIR = ROOT / "skills" / "pptx" / "togaf-architecture" / "loaders"
SHARED_DIR  = ROOT / "skills" / "pptx" / "shared"

sys.path.insert(0, str(SCRIPTS_DIR))
sys.path.insert(0, str(LOADERS_DIR))
sys.path.insert(0, str(SHARED_DIR))

from pptx_lib import build_pptx, header, layout_by_names, verify_pptx
from registry import LOADER_REGISTRY
from renderer_utils import configure_theme

HERE = Path(__file__).resolve().parent
CONFIG_FILE = HERE / "data" / "config_template.yaml"
MASTER_LIB  = ROOT / "skills" / "pptx" / "master-library"

MASTERS = [
    ("dark-cloudwise-green",  "togaf-dark-green.pptx"),
    ("light-cloudwise-cyan",  "togaf-light-cyan.pptx"),
]


def load_config(path):
    with open(path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def sort_loader_names(names):
    order = {"ba": 1, "aa": 2, "da": 3, "ta": 4}
    def _key(name):
        m = re.match(r"(ba|aa|da|ta)_(\d+)_", name)
        return (order.get(m.group(1), 999), int(m.group(2))) if m else (999, 999)
    return sorted(names, key=_key)


def add_cover_slide(ctx, cfg):
    spec_layout = ctx.template_spec.get("layout", {})
    cover_names = spec_layout.get("names", {}).get("cover", ["标题幻灯片", "封面"])
    cover_idx   = spec_layout.get("indices", {}).get("cover", 0)
    layout = layout_by_names(ctx.prs, cover_names, cover_idx)
    slide  = ctx.prs.slides.add_slide(layout)
    header(slide, cfg.get("title", "TOGAF Architecture Full Demo"))
    ph = next((s for s in slide.shapes
               if getattr(s, "is_placeholder", False) and s.placeholder_format.idx == 1), None)
    if ph and ph.has_text_frame:
        tf = ph.text_frame; tf.clear()
        p = tf.paragraphs[0]
        p.text = cfg.get("subtitle", "Loader Registry + Data Driven")
        p.font.size = Pt(ctx.typography["subtitle"])
        p.font.color.rgb = RGBColor(*ctx.colors["gray"])


def add_chapter_slide(ctx, title_text):
    spec_layout = ctx.template_spec.get("layout", {})
    section_names = spec_layout.get("names", {}).get("section", ["节标题", "目录页"])
    section_idx   = spec_layout.get("indices", {}).get("section",
                        spec_layout.get("indices", {}).get("cover", 0))
    layout = layout_by_names(ctx.prs, section_names, section_idx)
    slide  = ctx.prs.slides.add_slide(layout)
    header(slide, title_text)


def add_back_slide(ctx, cfg):
    spec_layout = ctx.template_spec.get("layout", {})
    thanks_names = spec_layout.get("names", {}).get("thanks", ["末尾幻灯片", "感谢页"])
    thanks_idx   = spec_layout.get("indices", {}).get("thanks", 0)
    layout = layout_by_names(ctx.prs, thanks_names, thanks_idx)
    slide  = ctx.prs.slides.add_slide(layout)
    header(slide, cfg.get("back_title", "Thank You"))


def my_slides(ctx):
    configure_theme(ctx.colors, ctx.template_spec)
    ctx.scale_x = 1.0  # TOGAF handles own scaling via _set_scale()
    cfg = load_config(CONFIG_FILE)
    add_cover_slide(ctx, cfg.get("meta", {}))
    for section in cfg.get("sections", []):
        add_chapter_slide(ctx, section.get("title", "Architecture Section"))
        for loader_name in sort_loader_names(list(section.get("slides", {}).keys())):
            loader_fn = LOADER_REGISTRY.get(loader_name)
            if not loader_fn:
                print(f"  [Warn] {loader_name} 未注册"); continue
            loader_fn(ctx, section["slides"][loader_name])
    add_back_slide(ctx, cfg.get("meta", {}))


def check_overflow(pptx_path):
    prs = Presentation(str(pptx_path))
    w = prs.slide_width
    bad = [i for i, sl in enumerate(prs.slides, 1)
           if max((sh.left + sh.width for sh in sl.shapes), default=0) > w]
    return len(prs.slides), bad


if __name__ == "__main__":
    print("=" * 60)
    for master_name, out_name in MASTERS:
        tpl  = MASTER_LIB / master_name / "cloudwise-master.pptx"
        spec = MASTER_LIB / master_name / "cloudwise-spec.yaml"
        out  = HERE / out_name
        print(f"\n▶ {master_name}")
        build_pptx(tpl, out, my_slides, spec)
        verify_pptx(out, tpl)
        slides, overflow = check_overflow(out)
        print(f"  Right overflow: {len(overflow)}/{slides}")
        if overflow:
            print(f"  Overflow slides: {overflow}")
    print("\n✓ All done")
