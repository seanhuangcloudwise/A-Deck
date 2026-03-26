#!/usr/bin/env python3
"""TOGAF architecture full demo generator (registry + data-driven)."""

import re
import sys
from pathlib import Path

import yaml
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parent.parent.parent
SCRIPTS_DIR = ROOT / "skills" / "pptx" / "scripts"
LOADERS_DIR = ROOT / "skills" / "pptx" / "togaf-architecture" / "loaders"
SHARED_DIR = ROOT / "skills" / "pptx" / "shared"

sys.path.insert(0, str(SCRIPTS_DIR))
sys.path.insert(0, str(LOADERS_DIR))
sys.path.insert(0, str(SHARED_DIR))

from pptx_lib import build_pptx, header, layout_by_names, verify_pptx  # noqa: E402
from registry import LOADER_REGISTRY  # noqa: E402
from renderer_utils import configure_theme  # noqa: E402

TEMPLATE = (
    ROOT
    / "skills"
    / "pptx"
    / "master-library"
    / "light-cloudwise-purple"
    / "cloudwise-master.pptx"
)
TEMPLATE_SPEC = (
    ROOT
    / "skills"
    / "pptx"
    / "master-library"
    / "light-cloudwise-purple"
    / "cloudwise-spec.yaml"
)
OUTPUT = Path(__file__).resolve().parent / "togaf-architecture-full-demo.pptx"
CONFIG_FILE = Path(__file__).resolve().parent / "data" / "config_template.yaml"


def load_config(path):
    with open(path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def sort_loader_names(names):
    order = {"ba": 1, "aa": 2, "da": 3, "ta": 4}

    def _key(name):
        m = re.match(r"(ba|aa|da|ta)_(\d+)_", name)
        if not m:
            return (999, 999)
        return (order.get(m.group(1), 999), int(m.group(2)))

    return sorted(names, key=_key)


def add_cover_slide(ctx, cfg):
    spec_layout = ctx.template_spec.get("layout", {})
    cover_names = spec_layout.get("names", {}).get("cover", ["标题幻灯片", "封面"])
    cover_idx = spec_layout.get("indices", {}).get("cover", 0)
    layout = layout_by_names(ctx.prs, cover_names, cover_idx)
    slide = ctx.prs.slides.add_slide(layout)
    header(slide, cfg.get("title", "TOGAF Architecture Full Demo"))

    subtitle_ph = next(
        (
            s
            for s in slide.shapes
            if getattr(s, "is_placeholder", False)
            and s.placeholder_format.idx == 1
        ),
        None,
    )
    if subtitle_ph and subtitle_ph.has_text_frame:
        tf = subtitle_ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = cfg.get("subtitle", "Loader Registry + Data Driven")
        p.font.size = Pt(ctx.typography["subtitle"])
        p.font.color.rgb = RGBColor(*ctx.colors["gray"])


def add_chapter_slide(ctx, title_text):
    spec_layout = ctx.template_spec.get("layout", {})
    section_names = spec_layout.get("names", {}).get("section", ["节标题", "目录页"])
    section_idx = spec_layout.get("indices", {}).get(
        "section",
        spec_layout.get("indices", {}).get("cover", 0),
    )
    layout = layout_by_names(ctx.prs, section_names, section_idx)
    slide = ctx.prs.slides.add_slide(layout)
    header(slide, title_text)


def add_back_slide(ctx, cfg):
    spec_layout = ctx.template_spec.get("layout", {})
    thanks_names = spec_layout.get("names", {}).get("thanks", ["末尾幻灯片", "感谢页"])
    thanks_idx = spec_layout.get("indices", {}).get("thanks", 0)
    layout = layout_by_names(ctx.prs, thanks_names, thanks_idx)
    slide = ctx.prs.slides.add_slide(layout)
    header(slide, cfg.get("back_title", "Thank You"))


def my_slides(ctx):
    configure_theme(ctx.colors, ctx.template_spec)
    # TOGAF loaders handle coordinate scaling internally via _set_scale();
    # neutralize build_pptx's scale_x to avoid double-scaling.
    ctx.scale_x = 1.0

    cfg = load_config(CONFIG_FILE)

    add_cover_slide(ctx, cfg.get("meta", {}))

    sections = cfg.get("sections", [])
    for section in sections:
        section_title = section.get("title", "Architecture Section")
        slides = section.get("slides", {})
        add_chapter_slide(ctx, section_title)

        for loader_name in sort_loader_names(list(slides.keys())):
            loader_fn = LOADER_REGISTRY.get(loader_name)
            if not loader_fn:
                print(f"[Warn] {loader_name} 未注册")
                continue
            print(f"[Slide] 正在生成 {loader_name}...")
            loader_fn(ctx, slides[loader_name])
            print("  ✓ 成功")

    add_back_slide(ctx, cfg.get("meta", {}))


def check_right_overflow(pptx_path):
    prs = Presentation(str(pptx_path))
    slide_width = prs.slide_width
    overflow_slides = []
    for i, slide in enumerate(prs.slides, 1):
        rightmost = max((sh.left + sh.width for sh in slide.shapes), default=0)
        if rightmost > slide_width:
            overflow_slides.append(i)
    return {
        "slides": len(prs.slides),
        "overflow_slides": overflow_slides,
    }


if __name__ == "__main__":
    print("=" * 60)
    print("TOGAF Architecture 全量生成器")
    print("=" * 60)

    print(f"已注册 {len(LOADER_REGISTRY)} 个装载器")
    build_pptx(TEMPLATE, OUTPUT, my_slides, TEMPLATE_SPEC)
    print(f"✓ PPT已生成: {OUTPUT}")
    verify_pptx(OUTPUT, TEMPLATE)
    overflow = check_right_overflow(OUTPUT)
    print(f"  Right overflow: {len(overflow['overflow_slides'])}/{overflow['slides']}")
    if overflow["overflow_slides"]:
        print(f"  Overflow slides: {overflow['overflow_slides']}")
    print("✓ 验证完成")
