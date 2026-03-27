#!/usr/bin/env python3
"""Check PPTX overflow with focus on business-content shapes.

This report separates:
1) raw overflow: all shapes
2) business overflow: excluding likely master/layout decorative shapes
"""

from __future__ import annotations

import argparse
from pathlib import Path
from typing import Tuple

from pptx import Presentation


EMU_PER_INCH = 914400


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    try:
        return (shape.text_frame.text or "").strip()
    except Exception:
        return ""


def _is_overflow(shape, slide_w, slide_h) -> bool:
    l, t, w, h = shape.left, shape.top, shape.width, shape.height
    return l < 0 or t < 0 or (l + w) > slide_w or (t + h) > slide_h


def _is_likely_non_business(shape, slide_w, slide_h) -> bool:
    """Heuristic filter for master/layout decorative shapes.

    Rules are intentionally conservative:
    - placeholders are excluded
    - huge no-text shapes (>105% slide width/height) are excluded
    - very large no-text lines spanning >105% slide width are excluded
    """
    if getattr(shape, "is_placeholder", False):
        return True

    text = _shape_text(shape)
    l, t, w, h = shape.left, shape.top, shape.width, shape.height

    huge_shape = (w > int(slide_w * 1.05) or h > int(slide_h * 1.05))
    if huge_shape and not text:
        return True

    # shape_type 9 is line in python-pptx enum values rendered as string in logs.
    likely_line = str(getattr(shape, "shape_type", "")).endswith("(9)")
    if likely_line and not text and w > int(slide_w * 1.05):
        return True

    return False


def check_ppt(path: Path) -> Tuple[dict, list[str]]:
    prs = Presentation(str(path))
    slide_w, slide_h = prs.slide_width, prs.slide_height

    summary = {
        "file": str(path),
        "size_in": (round(slide_w / EMU_PER_INCH, 3), round(slide_h / EMU_PER_INCH, 3)),
        "slides": len(prs.slides),
        "raw_overflow": 0,
        "business_overflow": 0,
    }
    lines = []

    for idx, slide in enumerate(prs.slides, start=1):
        raw_cnt = 0
        biz_cnt = 0
        for shape in slide.shapes:
            if not _is_overflow(shape, slide_w, slide_h):
                continue
            raw_cnt += 1
            if not _is_likely_non_business(shape, slide_w, slide_h):
                biz_cnt += 1

        summary["raw_overflow"] += raw_cnt
        summary["business_overflow"] += biz_cnt
        if raw_cnt or biz_cnt:
            lines.append(f"slide {idx}: raw={raw_cnt}, business={biz_cnt}")

    return summary, lines


def main():
    parser = argparse.ArgumentParser(description="Check raw/business overflow in PPTX")
    parser.add_argument("pptx", type=Path, help="path to .pptx")
    args = parser.parse_args()

    summary, lines = check_ppt(args.pptx)

    print("=== BUSINESS OVERFLOW REPORT ===")
    print(f"file: {summary['file']}")
    print(f"size_in: {summary['size_in'][0]} x {summary['size_in'][1]}")
    print(f"slides: {summary['slides']}")
    print(f"raw_overflow_shapes: {summary['raw_overflow']}")
    print(f"business_overflow_shapes: {summary['business_overflow']}")
    if lines:
        print("-- per-slide --")
        for line in lines:
            print(line)


if __name__ == "__main__":
    main()
