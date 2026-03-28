#!/usr/bin/env python3
"""验证节点文字颜色动态适配"""

from pptx import Presentation

def get_luminance(rgb):
    """Calculate relative luminance per WCAG."""
    r, g, b = rgb
    return (r * 0.2126 + g * 0.7152 + b * 0.0722) / 255

def check_node_colors(pptx_path, slide_idx=2):
    """Check text vs fill colors on a sample slide."""
    prs = Presentation(pptx_path)
    if slide_idx >= len(prs.slides):
        return None
    
    slide = prs.slides[slide_idx]
    results = []
    
    for shape in slide.shapes:
        if not hasattr(shape, 'fill') or not hasattr(shape, 'text_frame'):
            continue
        if not shape.text_frame.text.strip():
            continue
        
        # Get fill color
        fill_color = None
        if shape.fill.type == 1:  # SOLID
            try:
                fill_color = tuple(shape.fill.fore_color.rgb)
            except:
                pass
        
        # Get text color
        text_color = None
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if hasattr(run.font.color, 'rgb'):
                    try:
                        text_color = tuple(run.font.color.rgb)
                    except:
                        pass
                    break
        
        if fill_color and text_color:
            fill_lum = get_luminance(fill_color)
            text_lum = get_luminance(text_color)
            contrast = abs(text_lum - fill_lum)
            text_type = "dark" if text_lum < 0.4 else "light"
            fill_type = "dark" if fill_lum < 0.4 else "light"
            
            status = "OK" if contrast > 0.2 else "LOW"
            
            results.append({
                "label": shape.text_frame.text[:20],
                "fill_type": fill_type,
                "text_type": text_type,
                "contrast": round(contrast, 3),
                "status": status
            })
    
    return results

# Test
for name in ['bpmn-dark-cloudwise-green', 'bpmn-light-cloudwise-cyan', 'bpmn-light-cloudwise-purple']:
    pptx_path = f'projects/bpmn-architecture-full-demo/{name}.pptx'
    print(f"\n{name}:")
    colors = check_node_colors(pptx_path, slide_idx=2)
    if colors:
        ok_count = sum(1 for c in colors if c['status'] == 'OK')
        print(f"  Sample nodes: {len(colors)} | Good contrast: {ok_count}")
        for item in colors[:3]:
            print(f"    [{item['status']}] {item['label']:<15} fill:{item['fill_type']:<5} text:{item['text_type']:<5} contrast:{item['contrast']}")
